#!/usr/bin/env python3
"""
Generate the DeepSeek V4 Pro FPGA Inference Cluster Proposal PPT.
Enhanced content density with professional formatting.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Color Palette ──────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x27, 0x44)  # Dark navy
ACCENT     = RGBColor(0xE8, 0x77, 0x22)  # Orange/amber
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY   = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
SUBTITLE   = RGBColor(0x66, 0x66, 0x66)
TABLE_HDR  = RGBColor(0x1A, 0x27, 0x44)
TABLE_ROW1 = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_ROW2 = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_OK   = RGBColor(0x27, 0xAE, 0x60)
RED_BAD    = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW_WARN = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ── Helper Functions ───────────────────────────────────────

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_bg_rect(slide, color=DARK_BG):
    """Add full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_rich_textbox(slide, left, top, width, height):
    """Add empty textbox for rich multi-paragraph content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, font_size=16, color=DARK_TEXT, bold=False,
             alignment=PP_ALIGN.LEFT, space_after=Pt(6), font_name="Microsoft YaHei",
             first=False):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p

def add_slide_number(slide, num):
    tf = add_textbox(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                     str(num), font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, section_num, section_title, subtitle_text=None):
    """Add a chapter divider slide."""
    add_bg_rect(slide, DARK_BG)
    add_accent_bar(slide, Inches(1.5), Inches(3.2), Inches(0.08), Inches(0.8), ACCENT)
    tf = add_textbox(slide, Inches(1.8), Inches(2.8), Inches(10), Inches(0.6),
                     f"CH{section_num}  {section_title}", font_size=36, color=WHITE, bold=True)
    if subtitle_text:
        add_textbox(slide, Inches(1.8), Inches(3.6), Inches(10), Inches(0.5),
                    subtitle_text, font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC))
    add_slide_number(slide, section_num * 100)

def add_content_slide(slide, title, slide_num):
    """Standard content slide with title bar."""
    # White bg
    add_bg_rect(slide, WHITE)
    # Title bar
    add_accent_bar(slide, 0, 0, prs.slide_width, Inches(0.9), DARK_BG)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                title, font_size=26, color=WHITE, bold=True)
    # Bottom accent line
    add_accent_bar(slide, 0, Inches(7.35), prs.slide_width, Inches(0.05), ACCENT)
    add_slide_number(slide, slide_num)

def add_table(slide, left, top, col_widths, headers, rows, font_size=12):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w,
                                         Inches(0.36 * n_rows))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW1 if i % 2 == 0 else TABLE_ROW2
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
    return table

def add_code_block(slide, left, top, width, height, text, font_size=11):
    """Add a code/monospace text block."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x28, 0x2C, 0x34)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_right = Pt(12)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xAB, 0xB2, 0xBF)
        p.font.name = "Consolas"
        p.space_after = Pt(1)
    return tf

def add_ascii_diagram(slide, left, top, width, height, text, font_size=9,
                       bg_color=RGBColor(0xF0, 0xF4, 0xF8)):
    """Add a boxed ASCII diagram."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = MID_GRAY
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_right = Pt(8)
    tf.margin_bottom = Pt(6)
    for i, line in enumerate(text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Consolas"
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size + 2)
    return tf

def add_bullet_box(slide, left, top, width, height, items, font_size=14,
                   title=None, title_size=18, bg_color=LIGHT_GRAY):
    """Add a rounded bullet list box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(12)
    tf.margin_right = Pt(16)
    tf.margin_bottom = Pt(12)

    idx = 0
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = DARK_BG
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(8)
        idx = 1

    for item in items:
        if idx == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(5)
        idx += 1
    return tf

# ══════════════════════════════════════════════════════════════
# SLIDE 0: TITLE SLIDE
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_bg_rect(slide, DARK_BG)
add_accent_bar(slide, Inches(1.5), Inches(2.6), Inches(0.08), Inches(0.9), ACCENT)
add_textbox(slide, Inches(1.8), Inches(2.2), Inches(10), Inches(1.0),
            "DeepSeek V4 Pro — FPGA 算力推理集群",
            font_size=40, color=WHITE, bold=True)
add_textbox(slide, Inches(1.8), Inches(3.5), Inches(10), Inches(0.6),
            "可行性论证 & 工程设计方案",
            font_size=24, color=RGBColor(0xCC, 0xCC, 0xCC))
add_textbox(slide, Inches(1.8), Inches(4.5), Inches(10), Inches(0.5),
            "32 × Intel Agilex 7 M FPGA | 4 节点 RDMA 集群 | fp4 原生推理 | MLA 硬件加速",
            font_size=14, color=RGBColor(0x88, 0x88, 0x88))
add_textbox(slide, Inches(1.8), Inches(5.3), Inches(10), Inches(0.4),
            "内部评审 | 2025.05", font_size=14, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════
# AGENDA SLIDE
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_content_slide(slide, "目录 AGENDA", 0)
chapters = [
    ("CH1", "背景与战略定位", "为什么必须做、行业格局"),
    ("CH2", "架构总览", "三层架构、单节点拓扑、FPGA 内部 RTL"),
    ("CH3", "关键技术参数", "DeepSeek V4 Pro config.json 核心数值"),
    ("CH4", "算力分配设计", "32 卡资源切分、HBM 布局、性能推演"),
    ("CH5", "RTL 核心模块", "fp4 脉动阵列、MLA 流水线、KV Cache 管理"),
    ("CH6", "组网与通信", "400GbE RDMA 拓扑、带宽预算、F-Tile Ethernet"),
    ("CH7", "平台与物理形态", "算力卡规格、服务器选型、跨代兼容"),
    ("CH8", "软件生态", "OpenAI API 兼容、VFIO 驱动、推理服务层"),
    ("CH9", "开发路线图", "5 阶段递进、Go/No-Go 决策门"),
    ("CH10", "成本与财务", "原型 ¥12M、量产路径、TCO 分析"),
    ("CH11", "竞争分析", "vs 华为 Ascend / 国产 GPU / NVIDIA"),
    ("CH12", "风险评估与对策", "8 类风险 × 应对策略"),
]
for i, (ch, title, desc) in enumerate(chapters):
    y = Inches(1.3 + i * 0.48)
    add_accent_bar(slide, Inches(0.8), y + Inches(0.05), Inches(0.06), Inches(0.3), ACCENT)
    add_textbox(slide, Inches(1.1), y, Inches(1.2), Inches(0.35),
                ch, font_size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(2.4), y, Inches(5), Inches(0.35),
                title, font_size=14, color=DARK_TEXT, bold=True)
    add_textbox(slide, Inches(7.5), y, Inches(5), Inches(0.35),
                desc, font_size=11, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════
# CH1: 背景与战略定位
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 1, "背景与战略定位", "为什么必须做 | 行业格局 | 战略价值")

# ── Slide 1.1: 问题 ──
slide = add_blank_slide()
add_content_slide(slide, "中国大模型的部署瓶颈 — 模型能力与硬件底座的脱节", 101)

diagram = """
  ┌──────────────────────────────────────────────────────────────────┐
  │                                                                  │
  │    模型能力                          硬件底座                     │
  │   ════════════                    ════════════════               │
  │                                                                  │
  │   DeepSeek V4 Pro            ┌─────────────────────────┐        │
  │   1.6T MoE  全球顶尖         │ NVIDIA H100/B200        │        │
  │        │                     │ → 美国出口管制, 买不到   │        │
  │        │                     └─────────────────────────┘        │
  │        │                                                         │
  │   全球市场需求               ┌─────────────────────────┐        │
  │   东南亚 / 中东 / 拉美       │ 华为 Ascend 910B/C      │        │
  │   非洲 / 欧洲                │ → SMIC 7nm 产能受限     │        │
  │        │                     │ → CoWoS 封装制裁        │        │
  │        │                     │ → 无法全球出口          │        │
  │        │                     └─────────────────────────┘        │
  │        │                                                         │
  │        │                     ┌─────────────────────────┐        │
  │        └──────→  脱节 ←──────│ 寒武纪/海光/壁仞/摩尔    │        │
  │                              │ → 软件栈不成熟          │        │
  │                              │ → 供应同样受限          │        │
  │                              └─────────────────────────┘        │
  │                                                                  │
  │  核心矛盾: 能做出最好的模型, 但不知道用什么硬件部署出去           │
  └──────────────────────────────────────────────────────────────────┘
"""
add_ascii_diagram(slide, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.8), diagram, font_size=10)

# ── Slide 1.2: 答案 ──
slide = add_blank_slide()
add_content_slide(slide, "FPGA 路径 — 唯一可全球部署的中国大模型推理硬件", 102)

left_box = add_rich_textbox(slide, Inches(0.4), Inches(1.2), Inches(7.2), Inches(5.8))
add_para(left_box, "Intel Agilex 7 M + HBM2e 供应链分布", font_size=20, color=DARK_BG, bold=True, first=True)
add_para(left_box, "", font_size=6)
items_left = [
    "FPGA 芯片: Intel 自有 Fab (美国 / 爱尔兰 / 以色列)",
    "            不依赖 TSMC / SMIC",
    "HBM2e 堆叠: SK Hynix / Samsung (韩国)",
    "            两家供应商, 不是独家垄断",
    "先进封装:   Intel EMIB (马来西亚 / 越南)",
    "            非 CoWoS, 不受先进封装设备管制",
    "PCB 制造:    中国大陆 / 台湾 (本土自主)",
    "",
    "→ 供应链分散在 4 个以上国家/地区",
    "→ 不受任何单一司法管辖区完全控制",
]
for item in items_left:
    add_para(left_box, item, font_size=13, color=DARK_TEXT)

right_box = add_rich_textbox(slide, Inches(7.8), Inches(1.2), Inches(5.2), Inches(5.8))
add_para(right_box, "关键豁免", font_size=20, color=DARK_BG, bold=True, first=True)
add_para(right_box, "", font_size=6)
items_right = [
    "✓ 不受 GPU 算力出口管制",
    "   (TPP 远低于 4800 阈值)",
    "",
    "✓ 不受先进封装制裁",
    "   (EMIB ≠ CoWoS)",
    "",
    "✓ 标准 PCIe 5.0 CEM 设备",
    "   全球兼容任何服务器",
    "",
    "✓ 部署地点不受限",
    "   可在中国、东南亚、中东、",
    "   拉美、欧洲、非洲部署",
]
for item in items_right:
    add_para(right_box, item, font_size=12, color=DARK_TEXT, bold="✓" in item)

# ── Slide 1.3: 战略定位 ──
slide = add_blank_slide()
add_content_slide(slide, "战略定位 — 中国大模型出海唯一硬件底座", 103)

diag = """
  ┌─────────────────────────────────────────────────────────────────────┐
  │                                                                     │
  │           NVIDIA GPU              华为 Ascend              本方案 FPGA
  │         ┌────────────┐        ┌────────────┐         ┌────────────┐
  │  中国    │  ✗ 管制     │        │  ✓ 可售     │         │  ✓ 可获取   │
  │         └────────────┘        └────────────┘         └────────────┘
  │         ┌────────────┐        ┌────────────┐         ┌────────────┐
  │  东南亚  │  △ 降级版   │        │  △ 有限出口 │         │  ✓ 可售     │
  │         └────────────┘        └────────────┘         └────────────┘
  │         ┌────────────┐        ┌────────────┐         ┌────────────┐
  │  中东    │  △ 管制     │        │  ✗ 出口困难 │         │  ✓ 可售     │
  │         └────────────┘        └────────────┘         └────────────┘
  │         ┌────────────┐        ┌────────────┐         ┌────────────┐
  │  拉美    │  ✓ 可售     │        │  ✗ 无渠道   │         │  ✓ 可售     │
  │         └────────────┘        └────────────┘         └────────────┘
  │         ┌────────────┐        ┌────────────┐         ┌────────────┐
  │  欧洲    │  ✓ 可售     │        │  ✗ 出口困难 │         │  ✓ 可售     │
  │         └────────────┘        └────────────┘         └────────────┘
  │                                                                     │
  │  结论: 本方案是唯一在 "中国可获取" + "全球可部署" 两个维度上双满分    │
  └─────────────────────────────────────────────────────────────────────┘
"""
add_ascii_diagram(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.8), diag, font_size=11)

add_bullet_box(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.8),
    ["不是 \"FPGA 替代 GPU\"，而是 \"中国大模型出海 → 唯一可用的全球部署硬件路径\""],
    font_size=15, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ── Slide 1.4: 商业闭环 ──
slide = add_blank_slide()
add_content_slide(slide, "商业闭环 — 从训练到全球推理部署", 104)

diag = """
    中国境内                                   全球部署
   ════════════                            ═══════════════

   DeepSeek V4 Pro                         海外推理服务 API
   训练 (GPU 集群)              ──→         东南亚 / 中东 / 欧洲
        │                                   拉美 / 非洲
        ▼                                        ▲
   fp4 Checkpoint                               │
   权重导出                                    │
        │                                      │
        ▼                                      │
   ┌──────────────────────┐               ┌────────────────────┐
   │  FPGA 算力卡          │   装箱发货     │  FPGA 推理集群      │
   │  深圳/上海 生产        │ ────────────→ │  海外数据中心部署    │
   │  烧录 bitstream      │               │  即插即用 (标准PCIe) │
   │  预加载 V4 Pro 权重   │               │  OpenAI API 兼容    │
   └──────────────────────┘               └────────────────────┘

   关键:
   • FPGA 卡作为标准 PCIe 设备进入任何国家, 不触发 GPU 出口管制
   • 当地 Dell/Supermicro/HP 经销商采购服务器 → 插卡 → 上线
   • 用户只看到 OpenAI API, 不知道后端是 FPGA
"""
add_ascii_diagram(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.8), diag, font_size=10)

# ══════════════════════════════════════════════════════════════
# CH2: 架构总览
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 2, "架构总览", "三层架构 | 单节点拓扑 | FPGA 内部 RTL 模块")

# ── Slide 2.1: 三层架构 ──
slide = add_blank_slide()
add_content_slide(slide, "整体三层架构 — 从用户 API 到 FPGA 硬件", 201)

diag = """
  ┌──────────────────────────────────────────────────────────────────────┐
  │ [用户接入层]                                                          │
  │   OpenAI REST API 兼容端点                                            │
  │   /v1/chat/completions  /v1/completions  /v1/models                  │
  │   → 任何 OpenAI client (Python SDK / LangChain / Dify) 零成本接入      │
  └────────────────────────────┬─────────────────────────────────────────┘
                               │ HTTP/1.1
  ┌────────────────────────────▼─────────────────────────────────────────┐
  │ [推理服务层 — x86 主控服务器]                                          │
  │   Token 编码 (HuggingFace tokenizer)                                  │
  │   采样控制 (top-p / top-k / temperature)                              │
  │   多会话并发调度 + KV Cache 分配                                       │
  │   推理命令编排 → 结果拼接                                              │
  │   流式输出 (SSE) + FastAPI HTTP Server                               │
  └────────────────────────────┬─────────────────────────────────────────┘
                               │ 400GbE RoCE v2 RDMA
  ┌────────────────────────────▼─────────────────────────────────────────┐
  │ [FPGA 算力集群 — 4 × 8 卡节点, 32 FPGA]                               │
  │                                                                      │
  │  ┌──────────┐    ┌──────────┐    ┌──────────┐    ┌──────────┐      │
  │  │ Node 0   │    │ Node 1   │    │ Node 2   │    │ Node 3   │      │
  │  │ 8 FPGA   │    │ 8 FPGA   │    │ 8 FPGA   │    │ 8 FPGA   │      │
  │  │ L0~14    │───→│ L15~29   │───→│ L30~44   │───→│ L45~60   │      │
  │  │ +Embed   │    │          │    │          │    │+lm_head  │      │
  │  └──────────┘    └──────────┘    └──────────┘    └──────────┘      │
  │                                                                      │
  └──────────────────────────────────────────────────────────────────────┘
"""
add_ascii_diagram(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.8), diag, font_size=10)

# ── Slide 2.2: 单节点拓扑 ──
slide = add_blank_slide()
add_content_slide(slide, "单节点内部拓扑 — 标准 8-GPU 服务器复用", 202)

diag = """
          标准 4U GPU 服务器 (Supermicro SYS-821GE-TNHR)
    ┌──────────────────────────────────────────────────────────────┐
    │                                                               │
    │    Dual Intel Xeon (Sapphire Rapids, 160 PCIe 5.0 Lanes)     │
    │    ┌──────────────────┐       ┌──────────────────┐           │
    │    │     CPU 0        │◄─UPI─►│     CPU 1        │           │
    │    │   80 PCIe 5.0    │ 10.4  │   80 PCIe 5.0    │           │
    │    └──┬───┬───┬───┬──┘ GT/s  └──┬───┬───┬───┬──┘           │
    │       x16 x16 x16 x16          x16 x16 x16 x16               │
    │       │   │   │   │            │   │   │   │                  │
    │    ┌──┴───┴───┴───┴──┐      ┌──┴───┴───┴───┴──┐            │
    │    │   FPGA ×4        │      │   FPGA ×4        │            │
    │    │   Slot 0~3       │      │   Slot 4~7       │            │
    │    │   FHFL 双槽卡    │      │   FHFL 双槽卡    │            │
    │    └────────┬─────────┘      └────────┬─────────┘            │
    │             │      PCIe P2P          │                       │
    │             └──────────┬─────────────┘                       │
    │                        │                                     │
    │              ┌─────────▼──────────┐                          │
    │              │  F-Tile 200GbE    │                          │
    │              │  硬核 MAC + PCS    │                          │
    │              │  板载 QSFP-DD Cage │                          │
    │              └─────────┬──────────┘                          │
    └────────────────────────┼────────────────────────────────────┘
                             │ 直连 400GbE Switch
                             ▼
"""
add_ascii_diagram(slide, Inches(0.3), Inches(1.1), Inches(8.5), Inches(6.0), diag, font_size=9)

add_bullet_box(slide, Inches(9.0), Inches(1.2), Inches(4.0), Inches(5.5), [
    "同 CPU 下 FPGA P2P:",
    "  直通, ~260ns 延迟",
    "  无需经过 CPU 内存",
    "",
    "跨 CPU FPGA P2P:",
    "  经 UPI 链路",
    "  ~500ns 延迟",
    "  分级 All-Reduce 优化",
    "",
    "网络连接 (F-Tile):",
    "  板载 QSFP-DD Cage",
    "  直连 400GbE Switch",
    "  零 PCIe 中转",
    "",
    "供电:",
    "  4×3000W PSU (2+2)",
    "  每卡 ≤150W AUX",
], font_size=12)

# ── Slide 2.3: FPGA 内部 RTL ──
slide = add_blank_slide()
add_content_slide(slide, "单 FPGA 内部 RTL 模块 — 完整推理管线硬化", 203)

diag = """
     外部接口: PCIe 5.0 x8 CEM 金手指
              │
  ┌───────────▼────────────────────────────────────────────────────────┐
  │  PCIe 5.0 EP Hard IP (Intel 硬核, 零 LUT)                          │
  │  + 自定义推理报文控制器 (TLP ↔ 512-bit 推理帧)                       │
  └───────────┬────────────────────────────────────────────────────────┘
              │
  ┌───────────▼──────┬──────────┬──────────┬──────────┬───────────────┐
  │                  │          │          │          │               │
  │  fp4 Systolic    │   MLA    │  RoPE    │ RMSNorm  │  MoE Router   │
  │  Array ×8        │  Attn    │  Hard    │  Hard    │  + Dispatch   │
  │  (9,375 DSPs)    │  Pipe    │  Unit    │  Unit    │  + All-to-All │
  │                  │  line    │          │          │               │
  ├──────────────────┴──────────┴──────────┴──────────┴───────────────┤
  │  KV Cache Manager    │  Chip2Chip Router   │  Shared Expert Unit  │
  │  硬件哈希寻址        │  512-bit 报文组帧    │  单专家 66M MAC      │
  │  sliding_window=128  │  信用点反压流控     │  SiLU 硬化           │
  └──────────┬────────────────────────────────────────────────────────┘
             │
  ┌──────────▼────────────────────────────────────────────────────────┐
  │              HBM2e 控制器 (Avalon-MM, 2,048-bit @ 920 GB/s)        │
  │  ┌────────────────────────┬───────────────────────────────────┐   │
  │  │  权重区 (≤24 GB)        │  运行区 (≤8 GB)                    │   │
  │  │  · 12 路由专家 fp4     │  · KV Cache FP8                  │   │
  │  │  · 1 共享专家 fp4      │  · 激活 Buffer                   │   │
  │  │  · Attention 权重 fp4  │  · ETH TX/RX Ring Buffer         │   │
  │  │  · Router 权重         │  · 微批次中间激活                │   │
  │  └────────────────────────┴───────────────────────────────────┘   │
  │                        32 GB HBM2e (2 Stack)                       │
  └────────────────────────────────────────────────────────────────────┘
"""
add_ascii_diagram(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(6.0), diag, font_size=9)

# ══════════════════════════════════════════════════════════════
# CH3: 关键技术参数
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 3, "关键技术参数", "config.json 核心数值 | MLA | fp4 | 计算量分解")

# ── Slide 3.1: 架构参数表 ──
slide = add_blank_slide()
add_content_slide(slide, "DeepSeek V4 Pro 架构参数 — 来源: 开源 config.json", 301)

add_table(slide, Inches(0.4), Inches(1.2),
    [Inches(3.5), Inches(3.5), Inches(5.5)],
    ["参数", "值", "对 FPGA 设计的影响"],
    [
        ["hidden_size", "7,168", "所有矩阵乘的宽度基准"],
        ["num_hidden_layers", "61", "流水线分 4 段 (15+15+15+16)"],
        ["n_routed_experts", "384", "32 卡均分 = 12 专家/卡 (完美整除)"],
        ["n_shared_experts", "1", "每卡冗余存储 (33 MB fp4)"],
        ["num_experts_per_tok", "6", "All-to-All 通信: 6×7KB per token per layer"],
        ["moe_intermediate_size", "3,072", "每专家 66M MAC, 33 MB fp4 权重"],
        ["num_attention_heads", "128", "8 卡均分 = 16 头/卡 (TP 完美整除)"],
        ["num_key_value_heads", "1 (MLA!)", "KV 压缩为 1 个 576 维 latent vector"],
        ["head_dim (nope+rope)", "512 (448+64)", "Decoupled RoPE — 仅 64 维走位置编码"],
        ["q_lora_rank", "1,536", "Q 从 7168 压缩到 1536 → 再展开到 128×512"],
        ["o_lora_rank", "1,024", "O 从 128×512 压缩到 1024 → 再展开到 7168"],
        ["expert_dtype", "fp4 (E2M1)", "4-bit 浮点, 国产 GPU 全都不支持"],
        ["quantization_config", "FP8 (E4M3)", "激活值精度, 动态 block-wise 量化"],
        ["vocab_size", "129,280", "Embedding 1.85 GB FP16 (仅 Node 0)"],
        ["max_position_embeddings", "1,048,576", "1M context (FPGA 硬件 KV Cache 管理)"],
        ["sliding_window", "128", "局部注意力窗口, 硬件自动滑动"],
    ], font_size=13)

# ── Slide 3.2: MLA 详解 ──
slide = add_blank_slide()
add_content_slide(slide, "MLA (Multi-head Latent Attention) — 56× KV Cache 压缩", 302)

add_bullet_box(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(2.8), [
    "标准 MHA (如 LLaMA):",
    "  Q/K/V: 各自 128 heads × 128 dim",
    "  KV Cache: 2 × 128 × 128 = 32 KB/token/layer (FP16)",
    "  每 token 每层存储 32 KB, 1M context = 1.8 TB!",
    "",
    "DeepSeek MLA:",
    "  KV 压缩为 1 个 latent vector:",
    "    c_KV: 512 dim (nope, 不参与 RoPE)",
    "    k_R:   64 dim (decoupled RoPE)",
    "  KV Cache: (512+64) × FP8 = 576 Bytes/token/layer",
    "  压缩比: 32 KB → 576 B = 56×!",
    "  1M context × 61 layers × 576B = 33.6 GB",
], font_size=12)

add_bullet_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.8), [
    "MLA 在 FPGA 上的硬化优势:",
    "",
    "  GPU 方案: 需要自定义 CUDA/CANN",
    "  kernel 实现 soft Q/KV 压缩",
    "  每层需要 5+ 个独立 kernel launch",
    "  + LoRA 解压 + decoupled RoPE",
    "  → 大量 kernel launch overhead",
    "",
    "  FPGA 方案: 全部硬连线流水线",
    "  Q 压缩 → KV 压缩 → Q·K^T",
    "  → Softmax → A·V → O 解压",
    "  6 级流水线, 每级 ~3-9 μs",
    "  零 kernel launch, 零 CPU 参与",
], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.5), [
    "MLA 对 FPGA 的直接影响:",
    "  • KV Cache 仅 576 B/token/layer, 不需要跨 TP 卡切分 — 每卡存储完整 KV 副本",
    "  • 没有 V 投影矩阵 — MLA 的 V 直接复用 c_KV latent, 省掉 7168×4096 的 V 投影 MACs",
    "  • Q 和 O 走 LoRA 压缩/解压 — Attention 权重减少 ~70%, HBM 加载时间降低",
    "  • Decoupled RoPE — 仅 64 维做旋转位置编码, nope 部分(448 维)不参与, RoPE 硬件面积极小",
], font_size=12)

# ── Slide 3.3: fp4 详解 ──
slide = add_blank_slide()
add_content_slide(slide, "fp4 (E2M1) — 为什么国产 GPU 都不支持, FPGA 却可以", 303)

add_table(slide, Inches(0.4), Inches(1.2),
    [Inches(2.5), Inches(4.5), Inches(5.5)],
    ["GPU/芯片", "支持的推理精度", "fp4 支持?"],
    [
        ["NVIDIA H100/B200", "INT4/INT8/FP8/FP16/BF16", "✗ 不支持"],
        ["华为 Ascend 910B/C", "INT4/INT8/FP16/BF16", "✗ 不支持"],
        ["寒武纪 Siyuan 590", "INT4/INT8/FP16", "✗ 不支持"],
        ["海光 DCU (CDNA)", "INT8/FP16/BF16", "✗ 不支持"],
        ["壁仞 BR100", "INT8/FP16/BF16", "✗ 不支持"],
        ["本方案 FPGA", "fp4×fp8 原生乘法器", "✓ 定制 DSP 实现"],
    ], font_size=13)

add_bullet_box(slide, Inches(0.4), Inches(4.0), Inches(6.0), Inches(3.0), [
    "fp4 (E2M1) 格式:",
    "  {sign, 1b exponent, 2b mantissa}",
    "  有效值: {±0.5, ±1.0, ±1.5,",
    "           ±2.0, ±3.0, ±4.0, ±6.0}",
    "  15 个有效非零值",
    "",
    "GPU 不支持 fp4 的原因:",
    "  Tensor Core 是 INT8/FP8 硬化,",
    "  不支持浮点 4-bit 格式",
    "  需要先解压为 FP8 → 再算",
    "  → 浪费 HBM 空间和带宽",
], font_size=12)

add_bullet_box(slide, Inches(6.8), Inches(4.0), Inches(6.0), Inches(3.0), [
    "FPGA fp4 乘法器实现:",
    "  查表预计算 + DSP INT 模式",
    "  fp4 只有 16 种值 → 1 个 BRAM",
    "  存储 16×INT8 预计算缩放因子",
    "  DSP58 在 18×19 模式下 2×MAC/cycle",
    "  450 MHz → 8.44 TMACs/s per FPGA",
    "",
    "HBM 空间节省:",
    "  384 专家 (fp4):   12.7 GB",
    "  vs FP8:            25.4 GB",
    "  vs FP16:           50.7 GB",
    "  → 单卡 32GB HBM 轻松容纳",
], font_size=12)

# ── Slide 3.4: 计算量分解 ──
slide = add_blank_slide()
add_content_slide(slide, "单 Token 计算量精确分解 — 全 61 层 ~37.4 GMACs", 304)

add_table(slide, Inches(0.4), Inches(1.2),
    [Inches(3.0), Inches(2.5), Inches(2.5), Inches(4.5)],
    ["模块", "每 Token 每层 MACs", "全 61 层 MACs", "备注"],
    [
        ["MLA Q 压缩", "11.01M (7168×1536)", "0.67B", "LoRA down-projection"],
        ["MLA KV 压缩", "4.13M (7168×576)", "0.25B", "latent + rope parts"],
        ["MLA Q·K^T", "29.88M (128×448×512)", "1.82B", "128 头全并行"],
        ["MLA Softmax + A·V", "29.36M", "1.79B", "V 复用 c_KV latent"],
        ["MLA O 解压", "74.45M (LoRA ×2级)", "4.54B", "o_groups=16"],
        ["MLA 小计", "148.83M", "9.08B", ""],
        ["MoE 路由 (Gate)", "~1M", "~0.06B", "hash + scoring"],
        ["MoE 专家 FFN (×6)", "396.3M", "24.17B", "6×66.05M per expert"],
        ["MoE 共享专家", "66.05M", "4.03B", "每 token 都计算"],
        ["合计", "~612M", "~37.4B", ""],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.0), [
    "DSP 计算耗时:    37.4 GMACs / 8.44 TMACs/s = 4.43 ms  (纯计算, 无 HBM/通信开销)",
    "HBM 权重加载:    ~6.1 GB / 920 GB/s = 6.63 ms  (每 token 需从 HBM 加载的权重总量)",
    "结论: 计算与 HBM 接近平衡 (HBM 是微弱瓶颈, 1.36×), 说明 Agilex 7 M 的算力/HBM 配比恰好适合大模型推理",
], font_size=13, bg_color=RGBColor(0xFF, 0xF8, 0xE1))

# ══════════════════════════════════════════════════════════════
# CH4: 算力分配设计
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 4, "算力分配设计", "32 卡资源切分 | HBM 布局 | 性能推演")

# ── Slide 4.1: 32 卡切分 ──
slide = add_blank_slide()
add_content_slide(slide, "32 卡资源切分 — 4 节点 × 8 卡, 所有维度完美整除", 401)

add_table(slide, Inches(0.4), Inches(1.2),
    [Inches(1.5), Inches(1.5), Inches(2.0), Inches(2.2), Inches(2.2), Inches(3.0)],
    ["节点", "FPGA", "层范围", "独家专家", "Atten Head 数", "特殊"],
    [
        ["Node 0", "00~07", "Layer 00~14", "Expert 000~095", "16 头/卡", "Token Embedding"],
        ["", "", "15 层", "12/卡", "(128/8)", "1,850 MB FP16"],
        ["Node 1", "08~15", "Layer 15~29", "Expert 096~191", "16 头/卡", "—"],
        ["", "", "15 层", "12/卡", "(128/8)", ""],
        ["Node 2", "16~23", "Layer 30~44", "Expert 192~287", "16 头/卡", "—"],
        ["", "", "15 层", "12/卡", "(128/8)", ""],
        ["Node 3", "24~31", "Layer 45~60", "Expert 288~383", "16 头/卡", "lm_head + MTP"],
        ["", "", "16 层+MTP", "12/卡", "(128/8)", "1,850 MB FP16"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.5), [
    "完美整除验证:",
    "  128 heads / 8 cards = 16 heads/card  ✓  (TP 切分, 无余数)",
    "  384 experts / 32 cards = 12 experts/card  ✓  (专家并行, 无余数)",
    "  61 layers / 4 nodes ≈ 15~16 layers/node  ✓  (流水线并行, 负载接近均衡)",
    "",
    "为什么 8 卡而不是 10 卡?",
    "  • 8 卡是 AI 加速器的工业收敛标准 (DGX/HGX/Atlas 都是 8 卡)",
    "  • 128/8=16 整除, 128/10=12.8 不整除 → 之前 3×10 方案每个 TP group 有 2 张卡负载不均",
    "  • 标准 GPU 服务器原生支持 8×FHFL 插槽 → 不需要定制背板",
    "  • 384/32=12 整除, 384/30=12.8 不整除 → 每卡专家数对齐, HBM 布局完全一致",
], font_size=12)

# ── Slide 4.2: HBM 布局 ──
slide = add_blank_slide()
add_content_slide(slide, "单卡 HBM 内存地图 — 32 GB 空间分配", 402)

add_table(slide, Inches(0.4), Inches(1.2),
    [Inches(2.5), Inches(3.5), Inches(3.0), Inches(3.5)],
    ["区域", "内容", "大小", "备注"],
    [
        ["权重区 (~24GB 预算)", "", "", ""],
        ["", "12 路由专家 (fp4)", "396 MB", "12 × 33 MB, 独占不迁移"],
        ["", "1 共享专家 (fp4)", "33 MB", "全局冗余, 每卡都存"],
        ["", "Attention Q (fp4)", "~55 MB", "15~16 层 × 3.7 MB/layer"],
        ["", "Attention KV (fp4)", "~31 MB", "15~16 层 × 2.1 MB/layer"],
        ["", "Attention O (fp4)", "~118 MB", "15~16 层 × 7.9 MB/layer"],
        ["", "Router 权重 + LN", "~20 MB", "hash tables + RMSNorm"],
        ["", "Embedding (Node 0)", "1,850 MB", "129,280 × 7,168 FP16"],
        ["", "lm_head (Node 3)", "1,850 MB", "7,168 × 129,280 FP16"],
        ["", "权重小计 (Node 1/2)", "~589 MB", "远 < 24 GB 预算 ✓"],
        ["运行区 (~8GB)", "", "", ""],
        ["", "KV Cache", "~2.4-4.7 GB", "256K~512K context × 16 layers"],
        ["", "激活 Buffer", "~2 GB", "微批次中间激活"],
        ["", "ETH Ring Buffer", "~0.5 GB", "RDMA TX/RX 描述符环"],
        ["", "余量", ">25 GB", "可用于热门专家多副本 / 扩展 context"],
    ], font_size=12)

# ── Slide 4.3: 性能推演 ──
slide = add_blank_slide()
add_content_slide(slide, "性能推演 — Decode 延迟与集群吞吐", 403)

add_table(slide, Inches(0.4), Inches(1.2),
    [Inches(2.0), Inches(3.5), Inches(3.5), Inches(3.5)],
    ["场景", "每层延迟", "单节点吞吐", "说明"],
    [
        ["Decode (B=1)", "~65 μs", "~960 tok/s", "单 token 流水线, HBM 权重加载占 55%"],
        ["Decode (B=4)", "~150 μs", "~435 tok/s", "权重一次加载复用 4 tokens"],
        ["Decode (B=8)", "~280 μs", "~468 tok/s", "HBM 和 DSP 同时接近饱和"],
        ["Prefill (B=32)", "~1,657 μs", "~317 tok/s", "DSP 利用率接近 100%"],
        ["Prefill (B=128)", "~6,300 μs", "~333 tok/s", "大 batch prefill, DSP 满负荷"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(3.8), Inches(6.0), Inches(3.2), [
    "4 节点集群稳态吞吐:",
    "  B=1: ~1,000 tok/s (受最慢节点限制)",
    "  B=4: ~450 tok/s",
    "",
    "延迟分解 (Decode B=1, 每层 65μs):",
    "  HBM 权重加载: 35μs (54%)",
    "  DSP 计算:      23μs (35%)",
    "  跨节点通信:     5μs  (8%)",
    "  其他 (RMSNorm等): 2μs (3%)",
    "  → HBM 带宽是瓶颈, 不是计算",
], font_size=12)

add_bullet_box(slide, Inches(6.8), Inches(3.8), Inches(6.0), Inches(3.2), [
    "支持能力:",
    "  单 session:",
    "    首 token (128K prefill): ~2-3 秒",
    "    后续: ~960 tok/s 稳态",
    "",
    "  10 并发 session:",
    "    每 session ~80-100 tok/s",
    "    KV Cache 总占用 ~4.7 GB (256K)",
    "",
    "  20 并发 session:",
    "    每 session ~40-50 tok/s",
    "    KV Cache 总占用 ~9.4 GB",
    "    → 需降至 128K context 或 FP4 量化",
], font_size=12)

# ══════════════════════════════════════════════════════════════
# CH5: RTL 核心模块 (abbreviated - key slides)
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 5, "RTL 核心模块", "fp4 脉动阵列 | MLA 流水线 | KV Cache 管理 | MoE 调度")

# ── Slide 5.1: 顶层模块 ──
slide = add_blank_slide()
add_content_slide(slide, "顶层模块划分 — 13 个主模块 + 6 个子模块", 501)

add_table(slide, Inches(0.3), Inches(1.1),
    [Inches(3.5), Inches(4.5), Inches(4.5)],
    ["模块", "功能", "关键实现"],
    [
        ["pcie_cxl_ep_wrapper", "PCIe 5.0 x8 EP + CXL 可选", "Intel Hard IP, 零 LUT"],
        ["inference_ctrl_fsm", "全局推理状态机, 流水线控制", "Prefill/Decode 模式切换"],
        ["mla_attention_pipeline", "MLA Attention 6 级流水线", "Q/KV/O 压缩 + QK^T + AV + Softmax"],
        ["fp4_systolic_array", "8 个 128×128 脉动阵列", "fp4×fp8 MAC, 9,375 DSPs"],
        ["moe_expert_core", "MoE 专家 FFN + 路由", "gate/up/down + SiLU + top-6"],
        ["rope_hardware_unit", "Decoupled RoPE 硬化", "仅 64 维 rope part"],
        ["rms_norm_unit", "RMSNorm 硬化 (eps=1e-6)", "1 cycle 完成"],
        ["kv_cache_manager", "KV Cache 硬件寻址", "{session,layer,seq} → HBM addr"],
        ["hbm_memory_controller", "HBM2e 2048-bit 控制", "Avalon-MM, 920 GB/s"],
        ["chip2chip_router", "片间通信引擎", "512-bit 帧 + All-to-All + 流控"],
        ["shared_expert_unit", "共享专家 FFN", "每 token 都计算"],
        ["token_embed_lut", "Token Embedding (Node 0)", "129,280×7168 FP16"],
        ["lm_head_unit", "lm_head 投影 (Node 3)", "7168×129,280 FP16"],
    ], font_size=11)

# ── Slide 5.2: fp4 脉动阵列 ──
slide = add_blank_slide()
add_content_slide(slide, "fp4 脉动阵列 — 8 × 128×128 Systolic Array", 502)

add_bullet_box(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(3.0), [
    "fp4×fp8 乘法器实现:",
    "  方案: 查表预计算 + DSP INT 模式",
    "  • fp4 只有 2^4=16 种可能值 (15 有效)",
    "  • 1 个 BRAM (36Kb) 存储 16×INT8 缩放因子",
    "  • DSP58 在 18×19 模式下做 2×INT8 MAC/cycle",
    "  • 运行时: fp4 weight → BRAM lookup → DSP MAC",
    "  • 1 cycle 完成查表, 1 cycle 完成乘累加",
    "",
    "阵列配置:",
    "  8 个 128×128 脉动阵列",
    "  每阵列: 16,384 个 MAC 单元/cycle",
    "  合计: 131,072 MAC/cycle",
    "  @450 MHz: 59.0 GMACs/s per array",
    "  8 阵列: 8.44 TMACs/s total",
], font_size=12)

add_bullet_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(3.0), [
    "专家权重复用模式:",
    "  微批次 B=32 tokens →",
    "  同一专家权重加载一次 →",
    "  32 tokens 顺序流过脉动阵列",
    "",
    "HBM → 片上缓存流水线:",
    "  专家权重 (33MB) HBM→BRAM preload:",
    "    33MB / 920GB/s = 36 μs",
    "  32 tokens 激活值流过:",
    "    32 × 66M MAC / 8.44T = 250 μs",
    "  → 权重加载可和上一个专家计算重叠",
    "",
    "资源消耗:",
    "  9,375 DSP × 2 MAC/DSP",
    "  ~200 LUTs per multiplier × 17,056",
    "  ≈ 3.4M LUTs (超出 2.7M LE?)",
    "  → 大部分乘法器逻辑在 DSP 内部",
    "  → BRAM 查表共享, 不是每 MAC 独立 LUT",
], font_size=12)

# ── Slide 5.3: MLA 流水线 ──
slide = add_blank_slide()
add_content_slide(slide, "MLA Attention 流水线 — 6 级流水, 每 token 每层 ~25μs", 503)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(1.8), Inches(2.5), Inches(3.5), Inches(4.5)],
    ["流水级", "计算内容", "MACs", "延迟/瓶颈"],
    [
        ["Stage 0", "Q 压缩: 7168×1536", "11.01M", "6.0 μs (HBM 权重加载)"],
        ["Stage 1", "KV 压缩: 7168×576", "4.13M", "2.3 μs (HBM 权重加载)"],
        ["Stage 2", "Q·K^T: 128 头并行", "29.88M", "3.8 μs (DSP 计算)"],
        ["Stage 3", "Online Safe Softmax", "-", "0.2 μs (硬化, FP32 acc)"],
        ["Stage 4", "A·V: nope against c_KV", "29.36M", "3.7 μs (DSP, V 复用 c_KV)"],
        ["Stage 5", "O 解压: LoRA ×2级", "74.45M", "9.3 μs (HBM + DSP)"],
        ["合计", "MLA Attention", "148.83M", "~25 μs (部分流水重叠)"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(4.0), Inches(12.5), Inches(3.0), [
    "MLA 硬化 vs GPU 方案的关键差异:",
    "  1. 没有 V 投影矩阵 — MLA 的 V 直接复用 c_KV latent vector, 省掉 7168×4096(=29.4M) MAC 和对应的 HBM 权重加载",
    "  2. Q 和 O 走 LoRA 压缩 — 两级低秩分解 (q_lora_rank=1536, o_lora_rank=1024), Attention 权重总量减少 ~70% vs 标准 MHA",
    "  3. Decoupled RoPE — 只有 64 维 rope part 参与旋转位置编码, nope=448 维不参与, RoPE 硬件单元面积 < 2K LUTs",
    "  4. 零 kernel launch overhead — GPU 每层 Attention 需要 5+ 个独立 CUDA kernel, 每个 ~5-10μs launch 开销; FPGA 硬连线流水线, 零 CPU 参与",
    "  5. 128 个 Attention 头全部并行计算 — 每头有独立的 QK^T/AV MAC 单元, 不需要像 GPU 那样分 tile 串行",
], font_size=12)

# ══════════════════════════════════════════════════════════════
# CH6: 组网与通信
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 6, "组网与通信", "400GbE RoCE v2 | 拓扑与带宽 | F-Tile Ethernet")

# ── Slide 6.1: 拓扑 ──
slide = add_blank_slide()
add_content_slide(slide, "四节点 RDMA 集群拓扑 — 星型 400GbE Switch", 601)

diag = """
                       400GbE ToR Switch
                  (Spectrum SN4600C / 华为 CE8800)
                 ┌───────────────────────────────┐
                 │  RoCE v2 | DCQCN | PFC + ECN  │
                 │  16×400GbE 端口                │
                 └──┬──────┬──────┬──────┬────────┘
                    │      │      │      │
         400GbE QSFP-DD (或 2×200GbE)
                    │      │      │      │
              ┌─────┘      │      │      └─────┐
              ▼            ▼      ▼            ▼
    ┌──────────────┐ ┌──────────────┐ ┌──────────────┐ ┌──────────────┐
    │   Node 0     │ │   Node 1     │ │   Node 2     │ │   Node 3     │
    │  8 FPGA      │ │  8 FPGA      │ │  8 FPGA      │ │  8 FPGA      │
    │  F-Tile 200GbE │ │  F-Tile 200GbE │ │  F-Tile 200GbE │ │  F-Tile 200GbE │
    │  Layer 0-14  │ │  Layer 15-29 │ │  Layer 30-44 │ │  Layer 45-60 │
    │  +Embedding  │ │              │ │              │ │  +lm_head    │
    └──────────────┘ └──────────────┘ └──────────────┘ └──────────────┘
         │                  │                  │                  │
         └──────────────────┴──────────────────┴──────────────────┘
                        通信流 (全部经 Switch):
                   ① 组内 TP All-Reduce (节点内 P2P)
                   ② 组内 MoE Dispatch  (节点内 P2P)
                   ③ 跨节点 MoE Dispatch (经 Switch RDMA)
                   ④ 流水线边界 (Node i → Node i+1)
"""
add_ascii_diagram(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(6.0), diag, font_size=9)

# ── Slide 6.2: 带宽预算 ──
slide = add_blank_slide()
add_content_slide(slide, "片间通信带宽预算 — 所有链路利用率 < 10%", 602)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(3.0), Inches(2.5), Inches(2.5), Inches(2.2), Inches(2.2)],
    ["通信模式", "路径", "每层数据量", "带宽需求(500tok/s)", "链路利用率"],
    [
        ["TP All-Reduce", "节点内 PCIe P2P", "~224 KB", "213 MB/s", "0.8% (x16=28GB/s)"],
        ["组内 MoE Dispatch", "节点内 PCIe P2P", "~640 KB", "620 MB/s", "2.2%"],
        ["跨节点 MoE Dispatch", "RDMA via Switch", "~640 KB", "830 MB/s", "3.3% (200GbE)"],
        ["Pipeline 边界", "RDMA via Switch", "~7 KB", "<1 MB/s", "~0.01%"],
        ["KV Cache 同步", "RDMA via Switch", "~0.3 KB", "<100 MB/s", "~0.4%"],
        ["合计 (单向)", "—", "—", "~2.7 GB/s", "9.6% (x8 PCIe)"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(4.0), Inches(12.5), Inches(3.0), [
    "带宽结论:",
    "  • 每 FPGA → Switch: x8 PCIe 5.0 提供 28 GB/s 有效带宽, 需求 2.7 GB/s → 利用率 9.6%",
    "  • 跨节点: 200GbE (25 GB/s) 需求 ~0.83 GB/s per node pair → 利用率 3.3%",
    "  • x16 完全可以降为 x8 以节省收发器资源 → 16 个收发器腾出来做扩展或调试",
    "  • 400GbE Switch 端口利用率 < 5%, 有充足余量应对 burst 和未来扩容",
    "",
    "延迟分析:",
    "  • 同 CPU FPGA P2P:       ~260 ns  (PCIe TLP 直通)",
    "  • 跨 CPU FPGA P2P (UPI): ~500 ns  (UPI 转发)",
    "  • 跨节点 RDMA:            ~3 μs   (PCIe P2P + RoCE + Switch + PCIe P2P)",
    "  • MoE 全对全调度 58 层:   ~87 μs  (仅占单 token 总延迟 4ms 的 2.2%)",
], font_size=12)

# ── Slide 6.3: F-Tile Ethernet ──
slide = add_blank_slide()
add_content_slide(slide, "F-Tile 内置 Ethernet — 无需外部 NIC", 603)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(3.0), Inches(9.5)],
    ["F-Tile 特性", "规格"],
    [
        ["Ethernet MAC 硬核", "400GbE / 4×100GbE / 2×200GbE (IEEE 802.3bs)"],
        ["RS-FEC", "硬核 (IEEE 802.3 Clause 134), 零 LUT 开销"],
        ["PCS", "硬核 100G/200G/400G"],
        ["收发器", "G8 SerDes up to 116 Gbps"],
        ["连接方式", "板载 QSFP-DD Cage → 直连 400GbE Switch"],
        ["RoCE v2", "FPGA 软逻辑实现 (RDMA 协议栈)"],
        ["PCIe 占用", "零 (不经 PCIe, 数据面直出交换机)"],
        ["vs 外部 NIC", "省 ¥12K/节点 + 1 PCIe 槽 + ~20W 功耗"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.7), [
    "F-Tile 方案的核心优势:",
    "  1. 零 PCIe 中转: FPGA 推理数据直接封装 Ethernet Frame → 交换机, 零 CPU 参与数据面",
    "  2. 延迟降低: 省掉 FPGA → NIC PCIe DMA (~1μs), 数据从 HBM 直出 F-Tile MAC",
    "  3. 成本节省: 无需外部 NIC 芯片 (¥12K/ea × 4 = ¥48K), 无需备件 NIC",
    "  4. 功耗节省: F-Tile Ethernet 硬核 ~5W vs Intel E830 NIC ~25W",
    "  5. 供应简化: 不依赖任何 NIC 供应商, FPGA 卡自包含网络能力",
    "  6. 兼容模式: 仍可通过 PCIe 插标准 NIC (如华为 SP680) 用于混合网络场景",
    "",
    "数据路径对比:",
    "  旧 (NIC):  FPGA HBM → PCIe DMA → NIC → Ethernet → Switch  (+1μs, +25W)",
    "  新 (F-Tile): FPGA HBM → F-Tile MAC → Ethernet → Switch        (零中转)",
], font_size=11)

# ══════════════════════════════════════════════════════════════
# CH7: 平台与物理形态 (abbreviated)
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 7, "平台与物理形态", "算力卡规格 | 服务器选型 | 跨代兼容")

slide = add_blank_slide()
add_content_slide(slide, "FPGA 算力卡与服务器平台 — 标准 PCIe, 解耦迭代", 701)

add_bullet_box(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(5.5), [
    "FPGA 算力卡物理规范:",
    "",
    "  形态: FHFL (Full Height Full Length)",
    "        111.15mm × 312mm",
    "        双槽宽 (PCIe bracket ×2)",
    "",
    "  接口: PCIe 5.0 x16 CEM 金手指",
    "        标准边沿连接器",
    "",
    "  供电: 75W (PCIe 插槽)",
    "        + 150W (2× 8-pin AUX)",
    "",
    "  散热: 被动散热片 + 服务器风道",
    "        GPU 服务器前→后强力风道",
    "",
    "  管理: SMBus (I2C) → BMC",
    "        IPMI 标准温度/功耗上报",
    "",
    "  TDP:  ~75W (板卡) + 外部",
], font_size=12)

add_bullet_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.5), [
    "服务器平台: Supermicro 821GE ★",
    "",
    "  10× PCIe 5.0 x16 FHFL 插槽",
    "  双路 Xeon SPR (160 Lanes)",
    "  4× 3000W PSU (2+2 冗余)",
    "  X13→X14→X15 同机箱持续迭代",
    "",
    "备选平台:",
    "  Inspur NF5688M7 (国产首选)",
    "  H3C R5500 G6 (国产次选)",
    "",
    "跨代兼容保证:",
    "  2025: X13 + Agilex 7 M → Gen5",
    "  2027: X14 + 同一张卡 → Gen5 仍可用",
    "  2029: X15 + Agilex 10 M → Gen6",
    "  • 旧卡插新机: 降速 Gen5, 正常工作",
    "  • 新卡插旧机: 降速 Gen5, 正常工作",
    "  • 算力卡和平台完全解耦, 各自独立进化",
], font_size=12)

# ══════════════════════════════════════════════════════════════
# CH8: 软件生态 (abbreviated)
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 8, "软件生态", "OpenAI API 兼容 | VFIO 驱动 | 推理服务层")

slide = add_blank_slide()
add_content_slide(slide, "软件生态分层 — 用户只看到 OpenAI API", 801)

diag = """
  ┌─────────────────────────────────────────────────────────────┐
  │  应用开发者                                                   │
  │  import openai                                              │
  │  client = OpenAI(base_url="http://fpga-cluster:8080/v1")    │
  │  client.chat.completions.create(model="deepseek-v4", ...)   │
  │  → 与调用 GPT-4 API 完全一致, 不知道后端是 FPGA               │
  └──────────────────────────┬──────────────────────────────────┘
                             │ OpenAI REST API (HTTP/1.1 + SSE)
  ┌──────────────────────────▼──────────────────────────────────┐
  │  推理服务层 (x86 主控, 自研轻量调度器)                        │
  │  • Tokenizer (HuggingFace tokenizer)                        │
  │  • 采样器 (top-p / top-k / temperature)                     │
  │  • 多 session 并发管理 + KV Cache 分配                       │
  │  • 流式输出 (SSE) + FastAPI HTTP Server                    │
  └──────────────────────────┬──────────────────────────────────┘
                             │ libfpga.so (C 用户态库)
  ┌──────────────────────────▼──────────────────────────────────┐
  │  Linux VFIO 驱动 (无内核模块!)                                │
  │  /dev/vfio/N → mmap BAR2 → 直接读写 FPGA HBM                │
  │  MSI-X 中断通知推理完成                                       │
  │  IOMMU DMA 地址隔离                                          │
  └──────────────────────────┬──────────────────────────────────┘
                             │ PCIe 5.0 x8
  ┌──────────────────────────▼──────────────────────────────────┐
  │  FPGA 算力卡 (硬件, 固化 RTL)                                 │
  │  fp4 脉动阵列 + MLA + KV Cache + MoE Router                │
  │  32 GB HBM2e                                               │
  └─────────────────────────────────────────────────────────────┘
"""
add_ascii_diagram(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.2), diag, font_size=10)

add_bullet_box(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.9), [
    "兼容: OpenAI Python SDK / LangChain / LlamaIndex / Dify / FastGPT / Open WebUI / Continue.dev → 全部零成本接入",
    "不兼容 (也不需要): PyTorch Runtime / CUDA / HuggingFace Transformers → 模型在 GPU 训练后导出 fp4 checkpoint, 推理不经过这些框架",
], font_size=12, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ══════════════════════════════════════════════════════════════
# CH9: 开发路线图
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 9, "开发路线图", "5 阶段递进 | Go/No-Go 决策门 | 10 个月到全系统")

slide = add_blank_slide()
add_content_slide(slide, "五阶段递进开发 — 10 个月从单卡到 32 卡集群", 901)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(1.8), Inches(2.5), Inches(3.5), Inches(4.5)],
    ["阶段", "周期", "里程碑", "交付物"],
    [
        ["Phase 1", "Month 1-2", "单卡 fp4 精度验证", "HBM 带宽测试 + fp4 bit-accurate 对比"],
        ["单卡验证", "", "PCIe 5.0 链路 + HBM 读写", "单层推理 Micro-benchmark"],
        ["Phase 2", "Month 3-4", "8 卡 PCIe P2P 全互联", "TP All-Reduce + 组内 MoE"],
        ["8 卡节点", "", "15 层完整推理跑通", "吞吐 >200 tok/s target"],
        ["Phase 3", "Month 5-6", "F-Tile 200GbE + Switch", "RDMA 跨节点 FPGA→FPGA 通信"],
        ["双节点", "", "30 层流水线跑通", "跨节点 MoE Dispatch + Combine"],
        ["Phase 4", "Month 7-8", "32 卡全集群 61 层", "128K context + 多 session"],
        ["全集群", "", "全系统 Benchmark", "吞吐 >500 tok/s target"],
        ["Phase 5", "Month 9-10", "512K→1M context 极限", "热门专家 Multi-replica 优化"],
        ["生产优化", "", "故障注入 + Failover", "OpenAI API 兼容认证"],
    ], font_size=12)

# ── Slide 9.2: Phase 1 Go/No-Go ──
slide = add_blank_slide()
add_content_slide(slide, "Phase 1 单卡验证 — 决定项目去留的关键阶段", 902)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(4.5), Inches(4.0), Inches(4.0)],
    ["验证项", "通过标准", "不通过对策"],
    [
        ["HBM2e 读写带宽", "≥ 736 GB/s (80% 理论值)", "停项目"],
        ["fp4×fp8 乘法器精度", "单层 diff < 1e-3 vs PyTorch ref", "调试乘法器实现"],
        ["fp4 61 层累积精度", "diff < 2% vs PyTorch reference", "启动备选: 专家权重 fp8"],
        ["单层延迟", "< 200 μs (Attn+MoE)", "分析瓶颈, 优化 RTL"],
        ["PCIe 5.0 EP 链路", "Gen5 ×8 稳定, 无 CRC 错误", "降速 Gen4 调试"],
        ["FPGA 功耗", "< 85W (TDP 内)", "优化 DSP 利用率"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.2), [
    "Go/No-Go 硬停止条件 (任一触发, 项目暂停重评):",
    "  1. HBM 带宽实测 < 50% 理论值 (460 GB/s) → 停 (说明 HBM 控制器或板级设计有严重问题)",
    "  2. fp4 累积精度差 > 2% → 启动备选方案 (fp8 专家权重), 不直接停但需重评时间线和成本",
    "  3. Intel 确认供货周期 > 26 周 → 重新评估项目时间线",
    "",
    "Phase 1 成功标准: 以上 6 项全部绿灯 → 绿灯进入 Phase 2。",
    "Phase 1 预算: 2 张 Agilex 7 M 开发板 + 2 人月 RTL 开发 + 1 人月软件。",
], font_size=13, bg_color=RGBColor(0xFF, 0xF8, 0xE1))

# ══════════════════════════════════════════════════════════════
# CH10: 成本与财务
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 10, "成本与财务", "原型 ¥12M | 量产路径 | TCO 分析")

slide = add_blank_slide()
add_content_slide(slide, "原型开发预算 — 硬件 + 人力 + 其他", 1001)

add_table(slide, Inches(0.3), Inches(1.1),
    [Inches(3.5), Inches(3.5), Inches(2.5), Inches(3.0)],
    ["项目", "明细", "金额 (¥)", "备注"],
    [
        ["FPGA 芯片", "32×AGFB027 32GB HBM", "691,200", "¥21.6K/ea ($3K 实价)"],
        ["卡级物料 (1:1)", "PCB/VRM/散热/QSFP-DD/组装", "691,200", "¥21.6K/卡"],
                ["服务器机头", "4×Supermicro 821GE", "680,000", "¥170K/ea 裸机"],
        ["F-Tile 200GbE", "内置硬核 MAC (零外部芯片)", "0", "省 ¥12K/节点"],
        ["交换机", "1×400GbE (32×200GbE)", "100,000", "Spectrum SN4600C"],
        ["线缆/电源/机柜", "—", "60,000", ""],
        ["备件", "2 FPGA 卡", "86,400", "¥43.2K/卡"],
        ["硬件小计", "", "2,308,600", "≈ $321K"],
        ["", "", "", ""],
        ["FPGA RTL 人力", "5人×10月", "3,333,000", "¥800K/年"],
        ["软件/系统人力", "3人×10月", "1,500,000", "¥600K/年"],
        ["PCB 人力", "1人×5月", "250,000", "¥600K/年"],
        ["测试验证人力", "2人×8月", "667,000", "¥500K/年"],
        ["人力小计", "", "5,750,000", ""],
        ["", "", "", ""],
        ["其他 (工具/IP/FPGA SW/测试设备)", "", "2,000,000", ""],
        ["总计 (含 20% 余量)", "", "~12,000,000", "约 $1.7M"],
        ["", "", "", ""],
        ["物料 : 人工", "", "1 : 2.5", "硬件不贵, 人贵"],
    ], font_size=12)

# ── Slide 10.2: 量产路径 ──
slide = add_blank_slide()
add_content_slide(slide, "量产成本路径 — 从 ¥12M 原型到 ¥2M 量产", 1002)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(3.0), Inches(3.0), Inches(3.0), Inches(3.0)],
    ["", "原型 (1 套)", "小批量 (5 套)", "量产 (10+ 套)"],
    [
        ["单套硬件 BOM", "¥2.3M", "¥2.0M", "¥1.9M"],
        ["R&D NRE 摊薄", "¥9.7M/套", "¥1.9M/套", "¥0.97M/套"],
        ["单套总成本", "¥12M", "¥3.9M", "¥2.9M"],
        ["单卡物料成本", "¥43.2K", "¥37K", "¥33K"],
        ["", "", "", ""],
        ["FPGA 芯片价", "¥21.6K", "¥19K", "¥18K (小批量)"],
        ["外围物料", "¥21.6K (1:1)", "¥17K", "¥15K (批量 PCB/组装)"],
        ["服务器机头价", "¥170K", "¥160K", "¥150K (框架采购)"],
        ["", "", "", ""],
        ["物料 : 人工", "1 : 2.5", "1 : 0.7", "2.5 : 1"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.0), [
    "对比参考:",
    "  华为 Ascend 910B 8 卡服务器: ¥1.2-1.5M (但受产能限制, 大批量难以保证, 海外受限)",
    "  NVIDIA H100 8 卡服务器: ¥1.5-2M (中国无法正常获取, 出口管制)",
    "",
    "FPGA 核心优势: ¥33K/卡 + 可获取 + 可全球部署",
    "  量产 10 套后, 人工:物料 反转为 2.5:1 — 研发一次性摊完, 之后全是物料红利。",
], font_size=12, bg_color=RGBColor(0xFF, 0xF8, 0xE1))

# ══════════════════════════════════════════════════════════════
# CH11: 竞争分析
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 11, "竞争分析", "vs Ascend | vs 国产 GPU | vs NVIDIA | 护城河")

slide = add_blank_slide()
add_content_slide(slide, "对标矩阵 — FPGA 在关键维度上的唯一性", 1101)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)],
    ["维度", "NVIDIA B200", "华为 Ascend 910C", "国产 GPU(寒/海/壁)", "本方案 FPGA"],
    [
        ["中国可获取", "✗ 管制", "△ 排队", "△ 不定", "✓ 可获取"],
        ["全球可部署", "✓ 可售", "✗ 出口困难", "✗ 几乎无", "✓ 可部署"],
        ["供应稳定性", "✗ 切断", "△ 产能受限", "△ 不定", "✓ 多源供应"],
        ["fp4 原生支持", "✗", "✗", "✗", "✓ 定制"],
        ["MLA 硬件加速", "✗ 软件", "✗ 软件", "✗ 软件", "✓ 硬化"],
        ["软件生态", "★★★★★", "★★★★", "★★~★★★", "★★"],
        ["部署灵活性", "★★", "★★", "★★", "★★★★★"],
        ["$ / token", "最低(管制外)", "中", "中-高", "中"],
        ["运维成熟度", "★★★★★", "★★★★", "★★", "★★"],
    ], font_size=12)

add_bullet_box(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.0), [
    "结论: FPGA 方案在「可获取性」和「全球部署」两个维度上是唯一满分的方案。",
    "这构成了结构性优势, 其他方案无法通过技术迭代来复制。",
    "软件生态和运维成熟度是弱项, 需要通过自研推理服务层和标准化运维流程来弥补。",
], font_size=14, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ── Slide 11.2: FPGA 护城河 ──
slide = add_blank_slide()
add_content_slide(slide, "FPGA 护城河 — 三个不可复制的技术壁垒", 1102)

add_bullet_box(slide, Inches(0.4), Inches(1.2), Inches(3.8), Inches(5.5), [
    "护城河 ① fp4 原生推理",
    "",
    "所有 GPU (含 NVIDIA/AMD/",
    "Ascend/国产) → 不支持 fp4",
    "",
    "DeepSeek V4 Pro 权重 = fp4",
    "GPU 方案: fp4→解压→FP8→TC",
    "  → 浪费 HBM 空间 + 带宽",
    "",
    "FPGA: fp4 → fp4 MAC",
    "  → 零解压 → 全链路 fp4",
    "  → 4× HBM 节省 vs FP16",
    "  → 2× HBM 节省 vs FP8",
], font_size=12)

add_bullet_box(slide, Inches(4.5), Inches(1.2), Inches(4.2), Inches(5.5), [
    "护城河 ② MLA 硬件加速",
    "",
    "MLA = DeepSeek 独家 Attention",
    "GPU 方案: 需要定制 CUDA/CANN",
    "kernel 实现 Q/KV/O 压缩/解压",
    "每次 Attention 5+ kernel launch",
    "",
    "FPGA: 硬连线 MLA Pipeline",
    "  Q 压缩→KV 压缩→QK^T→",
    "  Softmax→AV→O 解压",
    "  6 级流水, 零 kernel launch",
    "  零 CPU 参与 Attention 计算",
], font_size=12)

add_bullet_box(slide, Inches(9.0), Inches(1.2), Inches(4.0), Inches(5.5), [
    "护城河 ③ KV Cache 硬件管理",
    "",
    "GPU 方案: vLLM PagedAttention",
    "  → Block Table 软件管理",
    "  → 每次 KV 访问走 CPU",
    "",
    "FPGA: 硬件哈希寻址",
    "  {session, layer, seq}",
    "  → HBM 物理地址",
    "  零 CPU 参与",
    "  Sliding Window (128)",
    "  硬件自动淘汰旧 token",
], font_size=12)

# ══════════════════════════════════════════════════════════════
# CH12: 风险评估
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 12, "风险评估与对策", "8 类风险 × 应对策略 | Go/No-Go")

slide = add_blank_slide()
add_content_slide(slide, "风险矩阵 — 8 类风险, 4 条硬停止条件", 1201)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(0.4), Inches(3.2), Inches(1.2), Inches(1.2), Inches(6.5)],
    ["#", "风险", "概率", "影响", "对策"],
    [
        ["1", "DeepSeek V5 改变架构 (放弃 MLA)", "中", "高",
         "跟踪V5研发动态; MLA已持续3代; 关键维度参数化"],
        ["2", "fp4 61层累积精度超标 >2%", "中", "高",
         "Phase1逐层bit-accurate验证; 备选fp8专家权重"],
        ["3", "FPGA被列入新一轮出口管制", "中低", "极高",
         "保持库存; 海外部署不受影响; 跟踪国产FPGA进展"],
        ["4", "Agilex 7 M 供货周期 >26周", "低", "高",
         "签供货协议; 备选Agilex 7 F-Series+外部DDR"],
        ["5", "PCIe P2P 跨CPU(UPI)兼容性", "中", "中",
         "Phase2早期验证; 备选插PCIe Switch卡统一"],
        ["6", "FPGA RTL 人才获取困难", "高", "中",
         "核心团队自建; 与高校FPGA实验室合作"],
        ["7", "Ascend突然支持fp4+MLA", "中", "高",
         "即使支持,供应仍受限; FPGA定位不改变"],
        ["8", "运维复杂度超预期", "中", "中",
         "BIST自检+远程JTAG+备卡热替换+详尽runbook"],
    ], font_size=11)

add_bullet_box(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.5), [
    "Go/No-Go 硬停止条件:",
    "  Phase 1 后: ① fp4 精度差异 >2% → 停/启用 fp8 备选 | ② HBM 带宽 <50% 理论值 → 停",
    "  Phase 4 后: ③ 32 卡吞吐 <300 tok/s → 评估经济性 | ④ 单卡月故障率 >1 → 重新设计散热/供电",
], font_size=12, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ══════════════════════════════════════════════════════════════
# CH13: 核心质疑回应 (评审修订 v1.1)
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_section_header(slide, 13, "核心质疑回应", "4 条大模型专家质疑 × 逐条论证 | v1.1 新增")

# ── Slide 13.1: 质疑 A — 实验验证计划 ──
slide = add_blank_slide()
add_content_slide(slide, "质疑 A: 没有实验数据？→ 开发板实证计划", 1301)

add_bullet_box(slide, Inches(0.3), Inches(1.1), Inches(6.2), Inches(6.0), [
    "三大关键实验 (Intel DK-SI-AGM027 ×2)",
    "",
    "实验 1 — fp4 精度验证 (最高优先级)",
    "  ① Python建模: 1层fp4量化 vs BF16",
    "  ② FPGA最小脉动阵列: 128×128",
    "  ③ 逐bit对比: DSP vs Python FP32",
    "  ④ 完整1层1000 token对比",
    "  判定: cosine similarity ≥0.995 ✓",
    "  <0.98 → 触发止损, 启用 fp8 备选",
], font_size=11)

add_bullet_box(slide, Inches(6.8), Inches(1.1), Inches(6.2), Inches(6.0), [
    "实验 2 — HBM 有效带宽实测",
    "  顺序1GB: ≥800 GB/s (近理论值)",
    "  MoE random access: power-law分布",
    "  12×33MB expert, α=1.2",
    "  判定: ≥550 GB/s (60%理论值)",
    "",
    "实验 3 — 单层端到端延迟",
    "  0-hit / 1-hit / 2-hit 三种情形",
    "  10,000 token统计长尾分布",
    "  判定: 加权平均 ≤15 μs/layer",
    "",
    "全部通过 → Phase 1 完成 → Phase 2",
], font_size=11)

add_bullet_box(slide, Inches(0.3), Inches(5.4), Inches(12.5), Inches(1.7), [
    "单卡可验证: fp4精度 / HBM带宽 / 单层延迟   |   多卡验证(Phase 2): 跨卡通信 / 全61层流水线 / 72h稳定性",
    "Go/No-Go: 实验1失败→评估fp8备选 | 实验2失败→重排weight layout | 实验3失败→重算TCO",
], font_size=11, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ── Slide 13.2: 质疑 B — HBM 带宽 ──
slide = add_blank_slide()
add_content_slide(slide, "质疑 B: 920 GB/s 不够？→ 带宽公平比较 + SRAM策略", 1302)

add_table(slide, Inches(0.3), Inches(1.2),
    [Inches(2.8), Inches(2.3), Inches(2.3), Inches(5.1)],
    ["指标", "H100 SXM", "FPGA Agilex 7M", "对比"],
    [
        ["权重精度", "BF16/FP16 (2B/param)", "fp4 (0.5B/param)", "4× 压缩"],
        ["HBM 原始带宽", "3.35 TB/s", "0.92 TB/s", "3.6× \"劣势\""],
        ["等效参数带宽", "1.68T params/s", "1.84T params/s", "FPGA +10%"],
        ["0-hit层 HBM读 (81.6%)", "~432 MB (全量)", "0 MB (SRAM常驻)", "FPGA 完胜"],
        ["1-hit层 HBM读 (16.9%)", "~432 MB", "33.4 MB", "FPGA 13× 少"],
        ["加权 HBM时间/层", "~129 μs", "7.2 μs", "FPGA 18× 快"],
        ["B=1 Decode瓶颈", "HBM带宽", "Expert单体大小", "不同瓶颈"],
        ["B=1 利用率", "~2-5%", "~49.5%", "FPGA 10-25×"],
    ], font_size=10)

add_bullet_box(slide, Inches(0.3), Inches(5.4), Inches(12.5), Inches(1.7), [
    "核心反论: fp4精度下等效参数带宽FPGA反超10%; SRAM缓存确定性权重(Shared+Attn+Router)消除81.6%层的全部HBM访问",
    "真正瓶颈不是920GB/s不够, 是Expert 33MB单体太大 — 即使HBM翻倍也无法完全消除 | 缓解: Expert预取/权重拆分",
], font_size=11, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ── Slide 13.3: 质疑 C — 软件栈 ──
slide = add_blank_slide()
add_content_slide(slide, "质疑 C: 软件栈从零开始？→ 85%复用, 自研仅14人月", 1303)

add_bullet_box(slide, Inches(0.3), Inches(1.1), Inches(6.2), Inches(3.8), [
    "复用开源 (零开发, ~85%代码等效)",
    "  Tokenizer: HuggingFace tokenizers",
    "  HTTP Server: FastAPI + uvicorn",
    "  Sampling: PyTorch/numpy (CPU端)",
    "  JSON Mode: LM Format Enforcer",
    "  Monitoring: Prometheus + Grafana",
    "  Auth/RateLimit: API Key + Redis",
    "  LangChain/Dify/Open WebUI: 原生HTTP",
], font_size=11)

add_bullet_box(slide, Inches(6.8), Inches(1.1), Inches(6.2), Inches(3.8), [
    "自研 (~14人月, ~15%代码量)",
    "  libfpga.so: VFIO/mmap/MMIO驱动",
    "  推理调度器: session/priority/round-robin",
    "  OpenAI API适配: protocol mapping",
    "  KV Cache管理器: 地址映射/prefix",
    "  权重加载器: PCIe DMA",
    "  → 总计 ~14,000行 C/Python",
    "",
    "GPU软件栈为什么庞大?",
    "  CUDA Driver→Runtime→cuBLAS→PyTorch→vLLM",
    "  FPGA: 算力在硬件里, 软件只管写寄存器",
], font_size=11)

add_bullet_box(slide, Inches(0.3), Inches(5.2), Inches(12.5), Inches(2.0), [
    "质疑中5个功能逐条: ①Continuous Batching→不需要(FPGA B=1已达50%利用) | ②Prefix Caching→有,硬件级 | ③SpecDec→v2 | ④P/D分离→天然支持",
    "三种使用者: ML开发者(OpenAI SDK,零成本) | 运维(Prometheus+Grafana,标准Linux) | FPGA开发者(我们5人,客户不需要碰Verilog)",
], font_size=11, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ── Slide 13.4: 质疑 D — 模型演进 ──
slide = add_blank_slide()
add_content_slide(slide, "质疑 D: V5变架构？→ DeepSeek第一性原理: 降低token成本", 1304)

add_bullet_box(slide, Inches(0.3), Inches(1.1), Inches(6.2), Inches(3.5), [
    "DeepSeek的技术三角 (缺一不可):",
    "  ① MoE: 算力成本降至Dense的~2%",
    "  ② fp4: 带宽需求降至BF16的1/4",
    "  ③ MLA: KV Cache压缩56×",
    "  三者互相强化, 放弃任一个成本跳升2-16×",
    "",
    "历史证据: V2→V3→V4 三代收敛",
    "  专家数: 160→256→384 (扩大)",
    "  MLA压缩比: 40×→50×→56× (增强)",
    "  V4引入fp4 (成本优化的自然延伸)",
    "  趋势: 更大模型+更低精度+更强压缩",
], font_size=11)

add_table(slide, Inches(6.8), Inches(1.1),
    [Inches(2.8), Inches(1.0), Inches(2.4)],
    ["V5可能变化", "概率", "RTL影响"],
    [
        ["专家数/层数/Top-K", "高", "寄存器改,零编译"],
        ["MLA维度/head数", "中", "寄存器改,零编译"],
        ["fp4→fp6 (提精度)", "低", "部分重配置~1h"],
        ["放弃 MLA", "极低", "重新设计,概率极低"],
        ["放弃 MoE (Dense)", "极低", "成本爆炸,不可能"],
        ["放弃 fp4 (fp8)", "极低", "成本翻倍,不可能"],
    ], font_size=10)

add_bullet_box(slide, Inches(0.3), Inches(5.0), Inches(12.5), Inches(2.2), [
    "高概率变化(参数调整)→仅改寄存器+WLC重生成(几分钟), 零编译 | 极低概率变化(放弃核心技术)→与DeepSeek使命矛盾",
    "退一步: 即使V5放弃MLA, 对GPU同样致命(KV Cache 576MB→32GB); V4 Pro市场需求窗口至少2-3年, 不依赖V5立即适配",
], font_size=11, bg_color=RGBColor(0xFF, 0xF0, 0xE0))

# ══════════════════════════════════════════════════════════════
# END SLIDE
# ══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_bg_rect(slide, DARK_BG)
add_accent_bar(slide, Inches(1.5), Inches(2.8), Inches(0.08), Inches(0.8), ACCENT)

add_textbox(slide, Inches(1.8), Inches(2.2), Inches(10), Inches(1.0),
            "战略层: 中国大模型出海 → 唯一可全球部署的推理硬件",
            font_size=22, color=RGBColor(0xCC, 0xCC, 0xCC))
add_textbox(slide, Inches(1.8), Inches(2.9), Inches(10), Inches(0.6),
            "────  ────  ────",
            font_size=18, color=MID_GRAY)
add_textbox(slide, Inches(1.8), Inches(3.3), Inches(10), Inches(0.6),
            "架构: 32 FPGA, 4×8 节点, 400GbE RDMA",
            font_size=18, color=WHITE)
add_textbox(slide, Inches(1.8), Inches(3.8), Inches(10), Inches(0.6),
            "技术: fp4 原生 + MLA 硬化 + HBM 常驻权重",
            font_size=18, color=WHITE)
add_textbox(slide, Inches(1.8), Inches(4.3), Inches(10), Inches(0.6),
            "生态: OpenAI API 兼容, Linux VFIO 标准驱动",
            font_size=18, color=WHITE)
add_textbox(slide, Inches(1.8), Inches(4.8), Inches(10), Inches(0.6),
            "财务: ¥12M 原型, 10 个月到全系统, ¥2M 量产",
            font_size=18, color=WHITE)
add_textbox(slide, Inches(1.8), Inches(5.3), Inches(10), Inches(0.6),
            "评审: v1.1 已完成 4 条核心质疑回应 (实验验证/HBM带宽/软件栈/模型演进)",
            font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB))

add_textbox(slide, Inches(1.8), Inches(6.4), Inches(10), Inches(0.6),
            "下一步: 采购 2 张 Agilex 7 M 开发板 → 三大实验验证 → 拿着 benchmark 找种子客户",
            font_size=16, color=ACCENT, bold=True)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = r"D:\workspace\fpgalpu\docs\fpga_inference_cluster_proposal.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
