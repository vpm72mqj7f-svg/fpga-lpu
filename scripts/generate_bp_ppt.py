#!/usr/bin/env python3
"""
Generate BP-style PPT for FPGA Inference Cluster project.

Audience:  generic (investor + strategic partner + internal exec)
Style:     data-driven, every slide has measurable numbers
Length:    14 slides

Data source: docs/e2e_validation_results.json (18 verified configs)
References:  fpga_inference_cluster_proposal.md (§4.6.1.7, §10.7, §11.6)

Output: docs/fpga_inference_bp.pptx
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ── Palette ──────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x27, 0x44)
ACCENT     = RGBColor(0xE8, 0x77, 0x22)
ACCENT2    = RGBColor(0x4A, 0x90, 0xE2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF7, 0xF9, 0xFC)
DARK_TEXT  = RGBColor(0x22, 0x33, 0x4A)
MID_TEXT   = RGBColor(0x55, 0x66, 0x80)
GREEN_OK   = RGBColor(0x27, 0xAE, 0x60)
RED_BAD    = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW     = RGBColor(0xF3, 0x9C, 0x12)
TABLE_HDR  = RGBColor(0x1A, 0x27, 0x44)
TABLE_ROW1 = RGBColor(0xF2, 0xF5, 0xFA)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────

def add_slide(bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    return slide


def add_box(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, l, t, w, h, text, size=18, color=DARK_TEXT, bold=False,
             align=PP_ALIGN.LEFT, font=FONT_CN, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font
    return tf


def add_paras(tf, paras):
    """paras: list of (text, size, color, bold, align) tuples"""
    for i, (text, size, color, bold, align) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = align
        p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = FONT_CN
        p.space_after = Pt(4)


def title_block(slide, num, title, subtitle=None):
    """Standard header: accent bar + slide number + title + subtitle"""
    # Accent bar (left edge)
    add_box(slide, Inches(0), Inches(0.4), Inches(0.18), Inches(1.0), ACCENT)
    # Slide number (top right)
    add_text(slide, Inches(12.4), Inches(0.3), Inches(0.8), Inches(0.3),
             f"{num:02d} / 21", size=10, color=MID_TEXT, align=PP_ALIGN.RIGHT)
    # Title
    add_text(slide, Inches(0.5), Inches(0.4), Inches(11.5), Inches(0.6),
             title, size=26, color=DARK_TEXT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.0), Inches(11.5), Inches(0.4),
                 subtitle, size=14, color=MID_TEXT)
    # Bottom rule
    add_box(slide, Inches(0.5), Inches(7.05), Inches(12.3), Emu(15000), ACCENT)


def metric_card(slide, l, t, w, h, value, label, sub=None,
                value_color=ACCENT, bg=LIGHT_BG):
    """Big number card."""
    add_box(slide, l, t, w, h, bg)
    add_text(slide, l, t + Inches(0.15), w, Inches(0.85),
             value, size=44, color=value_color, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(slide, l, t + Inches(0.95), w, Inches(0.35),
             label, size=13, color=DARK_TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, l, t + Inches(1.30), w, Inches(0.35),
                 sub, size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)


def add_table(slide, l, t, col_widths, headers, rows,
              hdr_bg=TABLE_HDR, hdr_color=WHITE,
              header_size=12, row_size=11, row_h=Inches(0.35)):
    rows_count = len(rows) + 1
    cols_count = len(headers)
    total_w = sum(col_widths)
    table_h = Inches(0.4) + row_h * len(rows)
    table = slide.shapes.add_table(rows_count, cols_count, l, t, total_w, table_h).table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    table.rows[0].height = Inches(0.4)
    for i in range(1, rows_count):
        table.rows[i].height = row_h
    # Header
    for j, h_text in enumerate(headers):
        c = table.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = hdr_bg
        tf = c.text_frame; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.text = h_text; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(header_size); p.font.color.rgb = hdr_color
        p.font.bold = True; p.font.name = FONT_CN
    # Rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            if i % 2 == 1:
                c.fill.solid(); c.fill.fore_color.rgb = TABLE_ROW1
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE
            tf = c.text_frame; tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
            p = tf.paragraphs[0]; p.text = str(val)
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            p.font.size = Pt(row_size); p.font.color.rgb = DARK_TEXT
            p.font.name = FONT_CN
    return table


# ── Load verified data ───────────────────────────────────────
with open(os.path.join(os.path.dirname(__file__), '..', 'docs',
                        'e2e_validation_results.json'), encoding='utf-8') as f:
    E2E = json.load(f)


# ═════════════════════════════════════════════════════════════
# Slide 1: Cover
# ═════════════════════════════════════════════════════════════
s = add_slide(DARK_BG)
# Top accent bar
add_box(s, 0, 0, prs.slide_width, Inches(0.15), ACCENT)
# Project name
add_text(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
         "DeepSeek V4 Pro 推理加速卡", size=20, color=ACCENT, bold=True)
add_text(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.5),
         "FPGA 集群验证 → ASIC 量产", size=48, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.5),
         "为中国大模型基础设施提供可获取的硬件路径",
         size=18, color=RGBColor(0xCC, 0xCC, 0xCC))
# Stats strip at bottom
add_box(s, 0, Inches(5.8), prs.slide_width, Inches(1.3), RGBColor(0x0F, 0x18, 0x2F))
strip_items = [
    ("$1.6T", "DeepSeek V4 模型规模"),
    ("32", "AGM 039 芯片 / 服务器"),
    ("5,790", "实测 tok/s (agent)"),
    ("$1.30", "/百万 token (100 套量产)"),
]
for i, (val, lbl) in enumerate(strip_items):
    x = Inches(0.5 + i * 3.1)
    add_text(s, x, Inches(5.95), Inches(2.9), Inches(0.55),
             val, size=28, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, x, Inches(6.5), Inches(2.9), Inches(0.35),
             lbl, size=11, color=WHITE, align=PP_ALIGN.CENTER)
# Footer
add_text(s, Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3),
         "项目商业方案 v1.0  |  内部保密  |  2026/05",
         size=10, color=MID_TEXT)



# ═════════════════════════════════════════════════════════════
# Slide 2: 术语速查 (Glossary)
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 2, "术语速查 — 后面所有页都会用到",
            "按 4 类分组. 后续若不熟悉术语, 可翻回本页对照")

# 4 columns, each with category header + 6-8 term lines
GLOSSARY_COLS = [
    ("硬件 / 芯片", [
        ("FPGA",     "现场可编程门阵列, 上板后可重配置数字逻辑"),
        ("ASIC",     "专用集成电路, 流片后逻辑固化, 性能/能效最优"),
        ("eASIC",    "结构化 ASIC, 仅金属层流片 (¥35M vs 全定制 ¥150M)"),
        ("HBM",      "高带宽显存, 堆叠 DRAM 直连芯片 (920 GB/s)"),
        ("DSP",      "数字信号处理单元, FPGA 内做 MAC 乘累加"),
        ("PCIe 5.0", "主板高速通道, x16 = 128 GB/s"),
        ("C2C",      "Chip-to-Chip, 卡内 FPGA 间 SerDes 直连"),
        ("AGM 039",  "Intel Agilex 7 M 系列芯片, 32GB HBM 集成"),
    ]),
    ("模型 / 算法", [
        ("fp4",      "4-bit 浮点权重 (16 个值), 比 fp16 体积小 4×"),
        ("fp8",      "8-bit 浮点激活, 训练后常用精度"),
        ("BF16",     "16-bit 脑浮点, 训练精度基线"),
        ("MoE",      "专家混合模型, 每 token 只激活子集 (V4 是 6/384)"),
        ("MLA",      "多头隐性注意力, KV 压缩 56× (DeepSeek 独家)"),
        ("KV Cache", "键值缓存, 上下文越长越大, 决定并发上限"),
        ("Top-K",    "MoE 每 token 路由的专家数 (DeepSeek V4: 6)"),
        ("QAT",      "量化感知训练, fp4 模型训练时即考虑精度损失"),
    ]),
    ("推理服务", [
        ("Prefill",  "输入提示词的批量处理阶段"),
        ("Decode",   "逐 token 输出阶段, 自回归"),
        ("TTFT",     "首 token 延迟, 用户感知"),
        ("TPOT",     "每 token 间隔 (生成速度)"),
        ("TPS",      "tokens / second, 吞吐"),
        ("SLA / P95", "服务等级目标 / 95 分位数 (95% 请求满足)"),
        ("Continuous Batching", "vLLM 招牌, 动态拼接多 session"),
        ("session",  "一次对话, 持有自己的 KV cache"),
    ]),
    ("工程 / 商业", [
        ("NRE",      "Non-Recurring Engineering, 一次性工程费 (流片)"),
        ("RTL",      "寄存器级硬件描述 (Verilog/SystemVerilog)"),
        ("IP",       "可复用硬件模块 (FPGA 阶段攒的 IP 可用于 ASIC)"),
        ("BOM",      "物料清单, 硬件总成本"),
        ("TCO",      "总拥有成本 (硬件 + 电费 + 运维, 通常按 3 年算)"),
        ("TAM",      "可服务总市场规模"),
        ("Quartus",  "Intel FPGA 综合 / 布局布线工具"),
        ("Signal Tap", "在线逻辑分析, 抓 RTL 内部信号无需停机"),
    ]),
]

# Layout: 4 columns × 8 rows. Each col 3.05" wide starting at x=0.4
COL_W = 3.05
START_X = 0.4
START_Y = 1.7

for col_idx, (cat, items) in enumerate(GLOSSARY_COLS):
    x = Inches(START_X + col_idx * COL_W)
    # Category header
    add_box(s, x, Inches(START_Y), Inches(COL_W - 0.15), Inches(0.35), ACCENT)
    add_text(s, x, Inches(START_Y), Inches(COL_W - 0.15), Inches(0.35),
             cat, size=12, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Term lines (size compressed to fit 8 per col)
    for row_idx, (term, defn) in enumerate(items):
        y = Inches(START_Y + 0.45 + row_idx * 0.6)
        add_text(s, x, y, Inches(1.05), Inches(0.55),
                 term, size=10, color=DARK_TEXT, bold=True, font=FONT_EN)
        add_text(s, x + Inches(1.05), y, Inches(COL_W - 1.25), Inches(0.55),
                 defn, size=8, color=MID_TEXT)


# ═════════════════════════════════════════════════════════════
# Slide 3: 关键对标参数的意义
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 3, "关键对标参数的意义",
            "后面的所有数字都围绕这 5 个问题, 先理解\"为什么这个数字重要\"")

QUESTIONS = [
    ("HBM 带宽 / 算力比",
     "为什么这个比值比单看算力更重要?",
     "LLM Decode 每生成 1 token 要遍历几乎全部权重 (~6 GB). 计算只用 37.4 GMACs.\n"
     "带宽不足 → 算力闲置. GPU 利用率仅 ~3%, FPGA 达 ~67%.\n"
     "→ 看 HBM/算力比, 才能判断硬件是否匹配 LLM decode 这个工作负载."),
    ("TTFT (首 token 延迟)",
     "为什么 P95 < 500ms 是 SLA 红线?",
     "用户感知: <300ms 即时, 300-1000ms 可接受, >1000ms 明显延迟.\n"
     "Chatbot 业务必须 P95 < 500ms, agent 场景可放宽到 1s.\n"
     "→ TTFT 决定用户是否觉得 \"卡顿\", 比平均吞吐更重要."),
    ("TPS / 吞吐",
     "聚合吞吐 vs 单 session 吞吐有什么区别?",
     "单 session (B=1): 一个用户感知的生成速度 (660-720 tok/s @ FPGA).\n"
     "聚合 (多 session): 多用户分摊的总产出 (5,800-8,500 tok/s @ FPGA).\n"
     "→ 单 session 决定用户体验, 聚合决定服务器 ROI."),
    ("$/百万 token",
     "为什么这是核心成本指标?",
     "DeepSeek V4 Pro 实际: ¥0.1/¥12/¥24 (缓存命中/未命中/输出), 混合 $1.46/M.\n"
     "FPGA 修正口径 $1.03-1.73/M, 100 套规模起优于 API.\n"
     "→ 自建在大流量私有部署经济性已胜出."),
    ("Accept Rate (接入率)",
     "为什么不只看吞吐, 还要看接入率?",
     "高负载下系统会饱和, 必须拒绝部分请求.\n"
     "Accept 88% 意味着 88% 到达请求被服务, 其余排队/拒绝.\n"
     "→ 比单看 TPS 更能反映可用性 (TPS 高但 accept 低 = 用户体验差)."),
]

# Layout: 5 rows, each row 1 question
for i, (param, ques, body) in enumerate(QUESTIONS):
    y = Inches(1.65 + i * 1.05)
    # Left: parameter name in colored badge
    add_box(s, Inches(0.5), y, Inches(2.6), Inches(0.95), ACCENT2)
    add_text(s, Inches(0.55), y + Inches(0.08), Inches(2.5), Inches(0.4),
             param, size=13, color=WHITE, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(0.55), y + Inches(0.5), Inches(2.5), Inches(0.4),
             ques, size=9, color=WHITE, align=PP_ALIGN.LEFT)
    # Right: explanation
    add_box(s, Inches(3.2), y, Inches(9.6), Inches(0.95), LIGHT_BG)
    tf = add_text(s, Inches(3.35), y + Inches(0.08), Inches(9.4), Inches(0.85),
                  "", size=10, color=DARK_TEXT)
    for j, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_CN
        p.space_after = Pt(1)


# ═════════════════════════════════════════════════════════════
# Slide 4: 问题
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 4, "中国大模型推理硬件的三重困境",
            "GPU 禁运 + 国产替代缺货 + fp4 没人原生支持")

# 3 cards
metric_card(s, Inches(0.5), Inches(1.7), Inches(4.0), Inches(1.8),
            "$0", "可获取的 H100/B200", "出口管制 EAR 3A090", value_color=RED_BAD)
metric_card(s, Inches(4.7), Inches(1.7), Inches(4.0), Inches(1.8),
            "12 月+", "Ascend 910C 排队", "SMIC 7nm 产能瓶颈", value_color=YELLOW)
metric_card(s, Inches(8.9), Inches(1.7), Inches(4.0), Inches(1.8),
            "0", "国产芯片支持 fp4 原生", "Ascend/寒武纪/海光均无", value_color=RED_BAD)

# Lower half: pain table
add_text(s, Inches(0.5), Inches(3.8), Inches(11.5), Inches(0.4),
         "客户实际处境:", size=16, color=DARK_TEXT, bold=True)
rows = [
    ("金融/医疗/政府 (隐私合规)", "数据不能出内网", "公共 API 不可用",          "✗ 业务上不了线"),
    ("中国大模型出海 (东南亚/中东)", "GPU 不能出口",       "Ascend 出口受限",           "✗ 海外部署无门"),
    ("有规模化推理需求公司",         "GPU 涨到 $25/hr",    "DeepSeek API 价格已被押注", "△ 长期受制于人"),
]
add_table(s, Inches(0.5), Inches(4.3),
          [Inches(3.8), Inches(2.8), Inches(2.8), Inches(2.1)],
          ["客户类型", "硬约束", "现有方案缺陷", "结果"],
          rows, header_size=12, row_size=11, row_h=Inches(0.5))

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
         "→ 我们的方案: FPGA 跨越供应链 + fp4 原生硬化, 同时满足 \"可获取 + 可部署 + 性价比\"",
         size=14, color=ACCENT, bold=True)


# ═════════════════════════════════════════════════════════════
# Slide 5: 方案
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 5, "方案: 单台 4U 服务器跑完 1.6T MoE 模型",
            "8 卡 × 4 颗 Intel Agilex 7 M = 32 芯片协同流水")

# Left: physical spec
add_text(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
         "物理形态", size=16, color=DARK_TEXT, bold=True)
specs = [
    ("服务器",    "1 台 Inspur NF5688M7 (4U, 5.3 kW)"),
    ("加速卡",    "8 张 × PCIe 5.0 x16 + 卡内 C2C 双环"),
    ("芯片",      "32 颗 Intel Agilex 7 M (AGM 039-F)"),
    ("HBM",       "32 × 32 GB = 1 TB 总容量, 聚合 29.4 TB/s"),
    ("算力",      "32 × 11.07 = 354 TMACS (fp4×fp8 原生)"),
    ("接口",      "标准 PCIe, 任意国家部署, 不受管制"),
]
for i, (k, v) in enumerate(specs):
    y = Inches(2.15 + i * 0.45)
    add_text(s, Inches(0.5), y, Inches(1.5), Inches(0.4),
             k, size=12, color=ACCENT, bold=True)
    add_text(s, Inches(2.0), y, Inches(4.5), Inches(0.4),
             v, size=12, color=DARK_TEXT)

# Right: how it scales
add_box(s, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.0), LIGHT_BG)
add_text(s, Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.4),
         "为什么 32 芯片够跑 1.6T?", size=16, color=DARK_TEXT, bold=True)
add_text(s, Inches(7.2), Inches(2.35), Inches(5.4), Inches(0.4),
         "fp4 权重压缩比 + MLA KV 56× 压缩",
         size=13, color=MID_TEXT, bold=True)

calc = [
    ("总权重 (fp4)",       "~800 GB"),
    ("每片 HBM 装载",     "~25 GB 权重 + ~7 GB KV"),
    ("KV per token (MLA)", "576 B  (vs 32 KB MHA, ×56)"),
    ("Layer / chip",       "29 × 2 layer + 3 × 1 layer = 61"),
    ("Expert / chip",      "384 / 32 = 12 (baseline) 或 14.3 (Hot Rep)"),
    ("Pipeline 深度",      "32 hop × ~50 μs = 1.5 ms / token"),
]
for i, (k, v) in enumerate(calc):
    y = Inches(2.95 + i * 0.5)
    add_text(s, Inches(7.3), y, Inches(2.5), Inches(0.4),
             k, size=11, color=MID_TEXT)
    add_text(s, Inches(9.8), y, Inches(2.9), Inches(0.4),
             v, size=11, color=DARK_TEXT, bold=True, font=FONT_EN)



# ═════════════════════════════════════════════════════════════
# Slide 6: 整体架构 - 服务器框图
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 6, "整体架构 (1/3): 4U 服务器框图",
            "Inspur NF5688M7 整机 = 8 卡 × 4 芯片 = 32 颗 AGM 039, 单机即完整集群")

# Outer chassis
chassis_x, chassis_y = Inches(0.5), Inches(1.7)
chassis_w, chassis_h = Inches(12.3), Inches(4.9)
add_box(s, chassis_x, chassis_y, chassis_w, chassis_h, RGBColor(0xEB, 0xEF, 0xF5))
# Chassis label
add_text(s, chassis_x + Inches(0.15), chassis_y + Inches(0.05),
         Inches(5), Inches(0.3),
         "Inspur NF5688M7  4U Server Chassis", size=10, color=MID_TEXT, bold=True)
add_text(s, chassis_x + Inches(7), chassis_y + Inches(0.05),
         Inches(5.1), Inches(0.3),
         "5.3 kW max  ·  风冷  ·  双路 Xeon 6", size=10, color=MID_TEXT,
         align=PP_ALIGN.RIGHT)

# Two CPU sockets (top row)
cpu_y = chassis_y + Inches(0.5)
for i, label in enumerate(["CPU 0  (Root Complex 0)", "CPU 1  (Root Complex 1)"]):
    x = chassis_x + Inches(0.4 + i * 5.95)
    add_box(s, x, cpu_y, Inches(5.65), Inches(0.7), RGBColor(0x4A, 0x90, 0xE2))
    add_text(s, x, cpu_y, Inches(5.65), Inches(0.7),
             label, size=13, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# PCIe 5.0 root labels
pcie_y = cpu_y + Inches(0.75)
add_text(s, chassis_x + Inches(0.4), pcie_y, Inches(5.65), Inches(0.25),
         "PCIe 5.0  x16 × 4 lanes", size=9, color=MID_TEXT,
         align=PP_ALIGN.CENTER, font=FONT_EN)
add_text(s, chassis_x + Inches(6.25), pcie_y, Inches(5.65), Inches(0.25),
         "PCIe 5.0  x16 × 4 lanes", size=9, color=MID_TEXT,
         align=PP_ALIGN.CENTER, font=FONT_EN)

# 8 FPGA cards (2 rows × 4 cards)
card_y0 = chassis_y + Inches(2.0)
card_w = Inches(2.7)
card_h = Inches(1.05)
card_gap_x = Inches(0.18)
for row in range(2):
    for col in range(4):
        idx = row * 4 + col
        x = chassis_x + Inches(0.4) + col * (card_w + card_gap_x)
        y = card_y0 + Inches(row * 1.25)
        add_box(s, x, y, card_w, card_h, RGBColor(0x1A, 0x27, 0x44))
        # Card title
        add_text(s, x, y + Inches(0.05), card_w, Inches(0.3),
                 f"Card {idx}", size=11, color=ACCENT, bold=True,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        # Chip content
        add_text(s, x, y + Inches(0.35), card_w, Inches(0.3),
                 "4 × AGM 039-F", size=10, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        # Layers (1 row layers each card carries)
        # 61 layers / 8 cards: 5 cards × 8 layers + 3 cards × 7 layers
        layers_per_card = [8, 8, 8, 8, 8, 7, 7, 7]
        start = sum(layers_per_card[:idx])
        end = start + layers_per_card[idx] - 1
        layer_range = f"L{start:02d}-{end:02d}"
        if idx == 0:
            layer_range += " + Emb"
        if idx == 7:
            layer_range += " + lm_head"
        add_text(s, x, y + Inches(0.65), card_w, Inches(0.3),
                 layer_range, size=9, color=RGBColor(0xCC, 0xCC, 0xCC),
                 align=PP_ALIGN.CENTER, font=FONT_EN)

# PCIe backplane connector strip
bp_y = card_y0 + Inches(2.5)
add_box(s, chassis_x + Inches(0.4), bp_y, Inches(11.5), Inches(0.35), ACCENT)
add_text(s, chassis_x + Inches(0.4), bp_y, Inches(11.5), Inches(0.35),
         "PCIe 5.0 Backplane  (P2P DMA, no ToR switch needed)",
         size=11, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom external interfaces
ext_y = chassis_y + chassis_h - Inches(0.3)
add_text(s, chassis_x + Inches(0.3), ext_y, Inches(11.5), Inches(0.3),
         "外部: 100GbE BMC  ·  双 PSU 2200W 80+ Titanium  ·  4U 标准机箱  ·  无外挂网络设备",
         size=9, color=MID_TEXT)

# Right-side annotations
ann_x = chassis_x + chassis_w + Inches(0.2)
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
         "关键点: 单机自成集群  |  卡间 P2P DMA 替代 ToR Switch  |  卡内 C2C 替代外部网络  |  电源/散热达标 4U 标准",
         size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# Slide 7: 整体架构 - 加速卡框图
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 7, "整体架构 (2/3): FPGA 加速卡框图",
            "单卡 4 颗 AGM 039 + 卡内 Chip-to-Chip 双环 + PCIe 金手指")

# PCB outline
pcb_x, pcb_y = Inches(1.0), Inches(1.7)
pcb_w, pcb_h = Inches(11.3), Inches(4.5)
add_box(s, pcb_x, pcb_y, pcb_w, pcb_h, RGBColor(0x0F, 0x44, 0x2A))  # PCB green
# PCB label
add_text(s, pcb_x + Inches(0.15), pcb_y + Inches(0.05), Inches(6), Inches(0.3),
         "FPGA Accelerator Card  (PCIe 5.0 full-height, full-length, 2-slot)",
         size=10, color=RGBColor(0xCC, 0xCC, 0xCC), bold=True)
add_text(s, pcb_x + Inches(7.5), pcb_y + Inches(0.05), Inches(3.7), Inches(0.3),
         "PCB 14-layer  ·  4-rail VRM  ·  ~550 W TDP",
         size=10, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.RIGHT)

# PCIe gold fingers (bottom edge)
gf_h = Inches(0.25)
gf_y = pcb_y + pcb_h - gf_h - Inches(0.05)
gf_x = pcb_x + Inches(2.0)
gf_w = Inches(4.8)
add_box(s, gf_x, gf_y, gf_w, gf_h, ACCENT)
add_text(s, gf_x, gf_y, gf_w, gf_h,
         "PCIe 5.0 x16 CEM Gold Fingers",
         size=10, color=DARK_TEXT, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)

# 4 FPGA chips arranged in 2x2
chip_w = Inches(2.4)
chip_h = Inches(2.1)
chip_x0 = pcb_x + Inches(0.6)
chip_y0 = pcb_y + Inches(0.55)
chip_gap_x = Inches(0.5)
chip_gap_y = Inches(0.25)
chip_positions = []  # (cx, cy) chip center for drawing C2C lines
for row in range(2):
    for col in range(2):
        idx = row * 2 + col
        x = chip_x0 + col * (chip_w + chip_gap_x)
        y = chip_y0 + row * (chip_h + chip_gap_y)
        # Chip body
        is_master = (idx == 0)
        body_color = ACCENT if is_master else RGBColor(0x33, 0x55, 0x88)
        add_box(s, x, y, chip_w, chip_h, body_color)
        # Chip title
        add_text(s, x, y + Inches(0.1), chip_w, Inches(0.35),
                 f"AGM 039-F  Chip {idx}", size=12, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        # Master indicator
        if is_master:
            add_text(s, x, y + Inches(0.42), chip_w, Inches(0.25),
                     "PCIe Master  (only this chip talks to host)",
                     size=8, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        # HBM stack indicator
        hbm_y = y + Inches(0.75)
        add_box(s, x + Inches(0.3), hbm_y, chip_w - Inches(0.6), Inches(0.4),
                RGBColor(0x88, 0xBB, 0xEE))
        add_text(s, x + Inches(0.3), hbm_y, chip_w - Inches(0.6), Inches(0.4),
                 "HBM2e  32 GB  ·  920 GB/s", size=9, color=DARK_TEXT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        # DSP/SRAM stats
        add_text(s, x + Inches(0.1), y + Inches(1.25), chip_w - Inches(0.2),
                 Inches(0.25), "12,300 DSP  ·  32 MB SRAM", size=9, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(s, x + Inches(0.1), y + Inches(1.55), chip_w - Inches(0.2),
                 Inches(0.25), "11.07 TMACs  ·  450 MHz", size=9, color=WHITE,
                 align=PP_ALIGN.CENTER, font=FONT_EN)
        # F-Tile SerDes ports
        add_text(s, x + Inches(0.1), y + Inches(1.85), chip_w - Inches(0.2),
                 Inches(0.2), "F-Tile SerDes ×2 → Dual Ring",
                 size=8, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)
        chip_positions.append((x + chip_w/2, y + chip_h/2))

# C2C dual ring lines: connect chip 0-1, 1-3, 3-2, 2-0 (clockwise Ring A)
from pptx.util import Emu as EmuClass
def add_line_connector(slide, x1, y1, x2, y2, color=ACCENT, width_pt=2):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn

# Draw C2C ring (curved-ish with straight segments)
# Use slight offsets to make the ring visible
c0, c1, c2, c3 = chip_positions  # [TL, TR, BL, BR] given row-major layout
# 0(TL) → 1(TR): top horizontal line
add_line_connector(s, c0[0], c0[1] - Inches(0.2), c1[0], c1[1] - Inches(0.2),
                   color=ACCENT, width_pt=2.5)
# 1(TR) → 3(BR): right vertical
add_line_connector(s, c1[0] + Inches(0.2), c1[1], c3[0] + Inches(0.2), c3[1],
                   color=ACCENT, width_pt=2.5)
# 3(BR) → 2(BL): bottom horizontal
add_line_connector(s, c3[0], c3[1] + Inches(0.2), c2[0], c2[1] + Inches(0.2),
                   color=ACCENT, width_pt=2.5)
# 2(BL) → 0(TL): left vertical
add_line_connector(s, c2[0] - Inches(0.2), c2[1], c0[0] - Inches(0.2), c0[1],
                   color=ACCENT, width_pt=2.5)

# Ring B (opposite direction) - dotted simulation with second color
add_line_connector(s, c0[0], c0[1] - Inches(0.1), c1[0], c1[1] - Inches(0.1),
                   color=ACCENT2, width_pt=1.5)
add_line_connector(s, c1[0] + Inches(0.1), c1[1], c3[0] + Inches(0.1), c3[1],
                   color=ACCENT2, width_pt=1.5)
add_line_connector(s, c3[0], c3[1] + Inches(0.1), c2[0], c2[1] + Inches(0.1),
                   color=ACCENT2, width_pt=1.5)
add_line_connector(s, c2[0] - Inches(0.1), c2[1], c0[0] - Inches(0.1), c0[1],
                   color=ACCENT2, width_pt=1.5)

# PCIe line from Chip 0 to gold fingers
add_line_connector(s, c0[0], c0[1] + Inches(1.0), gf_x + gf_w/2, gf_y,
                   color=ACCENT, width_pt=3)

# Right-side legend
legend_x = pcb_x + Inches(7.5)
legend_y = pcb_y + Inches(0.6)
add_box(s, legend_x, legend_y, Inches(3.5), Inches(2.5),
        RGBColor(0x22, 0x33, 0x4A))
add_text(s, legend_x, legend_y + Inches(0.1), Inches(3.5), Inches(0.3),
         "图例 / Legend", size=11, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER)
# Legend items
def add_legend_row(idx, color, text):
    y = legend_y + Inches(0.5 + idx * 0.35)
    add_box(s, legend_x + Inches(0.2), y, Inches(0.4), Inches(0.18), color)
    add_text(s, legend_x + Inches(0.7), y - Inches(0.02), Inches(2.8),
             Inches(0.25), text, size=9, color=WHITE)
add_legend_row(0, ACCENT,                     "C2C Ring A (顺时针)")
add_legend_row(1, ACCENT2,                    "C2C Ring B (逆时针)")
add_legend_row(2, ACCENT,                     "PCIe 5.0 x16 (Chip 0 only)")
add_legend_row(3, RGBColor(0x88, 0xBB, 0xEE), "HBM2e 堆叠 (EMIB 封装)")
add_legend_row(4, ACCENT,                     "PCIe 金手指 (CEM)")
add_legend_row(5, RGBColor(0x33, 0x55, 0x88), "Chip 1-3 (无 PCIe)")

# Bottom annotation
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
         "关键设计: 仅 Chip 0 有 PCIe Hard IP, Chip 1-3 通过 C2C 双环代理与 Host 通信 (省 R-Tile 成本)",
         size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# Slide 8: 整体架构 - FPGA 内部 RTL 框图
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 8, "整体架构 (3/3): FPGA 单芯片内部 RTL 框图",
            "13 个核心模块, 总计 54 人月自研 (含 8 人月集成余量)")

# Outer FPGA boundary
fpga_x, fpga_y = Inches(0.6), Inches(1.65)
fpga_w, fpga_h = Inches(12.1), Inches(5.0)
add_box(s, fpga_x, fpga_y, fpga_w, fpga_h, RGBColor(0xF7, 0xF9, 0xFC))
add_text(s, fpga_x + Inches(0.15), fpga_y + Inches(0.05),
         Inches(6), Inches(0.3),
         "AGM 039-F  FPGA Die  (12,300 DSP · 32 MB SRAM · 32 GB HBM2e)",
         size=10, color=MID_TEXT, bold=True)

# Top: PCIe + DMA Engine (Chip 0 only)
pcie_y = fpga_y + Inches(0.5)
add_box(s, fpga_x + Inches(0.3), pcie_y, Inches(5.5), Inches(0.55),
        RGBColor(0xE8, 0x77, 0x22))
add_text(s, fpga_x + Inches(0.3), pcie_y, Inches(5.5), Inches(0.55),
         "PCIe 5.0 EP Hard IP + DMA Engine  (Chip 0 only)",
         size=12, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# C2C SerDes (top right)
c2c_x = fpga_x + Inches(6.1)
add_box(s, c2c_x, pcie_y, Inches(5.7), Inches(0.55),
        RGBColor(0x4A, 0x90, 0xE2))
add_text(s, c2c_x, pcie_y, Inches(5.7), Inches(0.55),
         "C2C SerDes Dual Ring  (F-Tile × 2 lanes, all chips)",
         size=12, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Middle: 7 RTL modules in a row
mod_y = fpga_y + Inches(1.3)
mod_h = Inches(1.6)
mod_w = Inches(1.62)
mod_gap = Inches(0.07)
modules = [
    ("fp4 Systolic\nArray ×8", "12,300 DSP\n11 TMACs", RGBColor(0x27, 0xAE, 0x60), "10 人月"),
    ("MLA Attention\nPipeline", "Q/K/V 压缩\nRoPE 解耦", RGBColor(0x27, 0xAE, 0x60), "12 人月"),
    ("MoE Router\n+ Dispatch", "Top-6/384\n动态路由", RGBColor(0x4A, 0x90, 0xE2), "4 人月"),
    ("KV Cache\nManager", "硬件寻址\nPagedAttn", RGBColor(0x4A, 0x90, 0xE2), "6 人月"),
    ("RMSNorm\n+ RoPE Unit", "硬件 SiLU\n+ Rotary", RGBColor(0x88, 0xBB, 0xEE), "1.5 人月"),
    ("Embedding\n+ lm_head", "Token table\n+ MTP", RGBColor(0x88, 0xBB, 0xEE), "3 人月"),
    ("Inference\nControl FSM", "61 层流水\n调度", RGBColor(0xCC, 0x88, 0x33), "2 人月"),
]
for i, (name, detail, color, pm) in enumerate(modules):
    x = fpga_x + Inches(0.3) + i * (mod_w + mod_gap)
    add_box(s, x, mod_y, mod_w, mod_h, color)
    # Module name
    tf = add_text(s, x, mod_y + Inches(0.1), mod_w, Inches(0.5),
                  "", size=10, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER)
    for j, line in enumerate(name.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10); p.font.color.rgb = WHITE
        p.font.bold = True; p.font.name = FONT_CN
    # Detail
    tf2 = add_text(s, x, mod_y + Inches(0.75), mod_w, Inches(0.5),
                   "", size=8, color=WHITE, align=PP_ALIGN.CENTER)
    for j, line in enumerate(detail.split('\n')):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(8); p.font.color.rgb = WHITE
        p.font.name = FONT_EN
    # PM cost
    add_text(s, x, mod_y + mod_h - Inches(0.3), mod_w, Inches(0.25),
             pm, size=9, color=RGBColor(0xFF, 0xEE, 0xCC), bold=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)

# Bottom: HBM controller + HBM
hbm_y = fpga_y + Inches(3.1)
add_box(s, fpga_x + Inches(0.3), hbm_y, Inches(11.5), Inches(0.45),
        RGBColor(0x33, 0x55, 0x88))
add_text(s, fpga_x + Inches(0.3), hbm_y, Inches(11.5), Inches(0.45),
         "HBM2e Memory Controller  (Avalon-MM, 2048-bit interface)",
         size=12, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# HBM bank divided
hbm_bank_y = hbm_y + Inches(0.55)
hbm_bank_h = Inches(1.0)
# Weight region
add_box(s, fpga_x + Inches(0.3), hbm_bank_y, Inches(7.5), hbm_bank_h,
        RGBColor(0x88, 0xBB, 0xEE))
add_text(s, fpga_x + Inches(0.3), hbm_bank_y, Inches(7.5), Inches(0.3),
         "权重区  Weight Region  ~25 GB", size=11, color=DARK_TEXT, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, fpga_x + Inches(0.3), hbm_bank_y + Inches(0.35), Inches(7.5), Inches(0.6),
         "12 个路由专家 (fp4) · 1 共享专家 · Attention 权重 · Router 表",
         size=9, color=DARK_TEXT, align=PP_ALIGN.CENTER)
# KV region
add_box(s, fpga_x + Inches(7.9), hbm_bank_y, Inches(3.9), hbm_bank_h,
        RGBColor(0xEE, 0xBB, 0x88))
add_text(s, fpga_x + Inches(7.9), hbm_bank_y, Inches(3.9), Inches(0.3),
         "运行区  Runtime  ~7 GB", size=11, color=DARK_TEXT, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, fpga_x + Inches(7.9), hbm_bank_y + Inches(0.35), Inches(3.9), Inches(0.6),
         "KV Cache (FP8) · 激活 Buffer · C2C RX/TX Ring",
         size=9, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# Bottom annotation
add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
         "数据流: PCIe → DMA → SRAM 缓存 → DSP 阵列 ↔ HBM ↔ MLA/MoE 流水 → C2C 转发下一层",
         size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# Slide 9: 技术差异化 — HBM/算力比
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 9, "为什么 FPGA 跑 LLM Decode 比 GPU 更高效",
            "Decode 是 memory-bound, GPU 的算力大量闲置")

# Comparison table with HBM/算力 ratio
rows = [
    ("NVIDIA H100",  "3,350 GB/s",  "990 TFLOPS (FP8)",  "3.4",   "~3%",   ">$30K (管制)"),
    ("NVIDIA H200",  "4,800 GB/s",  "990 TFLOPS",        "4.8",   "~3%",   ">$40K (管制)"),
    ("NVIDIA B200",  "8,000 GB/s",  "2,250 TFLOPS",      "3.6",   "~5%",   ">$40K (管制)"),
    ("Ascend 910C",  "~2,000 GB/s", "~800 TFLOPS (FP8)", "2.5",   "~30%",  "排队 12 月"),
    ("AGM 039 (32×)", "29.4 TB/s 总",  "354 TMACS (fp4)",   "83",    "~67%",  "8-12 周交货"),
]
add_table(s, Inches(0.5), Inches(1.8),
          [Inches(2.4), Inches(2.0), Inches(2.4), Inches(1.8), Inches(1.8), Inches(2.0)],
          ["硬件", "HBM 带宽", "算力", "HBM/算力比", "Decode 利用率", "可获取性"],
          rows, header_size=11, row_size=11, row_h=Inches(0.5))

# Key insight box
add_box(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.7), LIGHT_BG)
add_text(s, Inches(0.8), Inches(5.25), Inches(11.7), Inches(0.4),
         "核心洞察: DeepSeek decode 需要 6 MAC/Byte 计算密度",
         size=15, color=DARK_TEXT, bold=True)
add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.4),
         "● H100/H200/B200: HBM/算力比 ~3-5, Tensor Core 97% 闲置等 HBM",
         size=12, color=MID_TEXT)
add_text(s, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.4),
         "● AGM 039: HBM/算力比 83, DSP 与 HBM 接近平衡, 利用率 67%",
         size=12, color=GREEN_OK)
add_text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
         "→ 不是 \"FPGA 比 GPU 快\", 是 \"对 LLM decode 这个特定负载, FPGA 架构上更匹配\"",
         size=13, color=ACCENT, bold=True)


# ═════════════════════════════════════════════════════════════
# Slide 10: fp4 + MLA 硬化
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 10, "三个 GPU 做不到的事 (国产芯片同样做不到)")

cards = [
    ("①", "fp4 原生计算",
     "GPU 路径: 加载 fp4 → 解压到 FP8 → Tensor Core MAC\nFPGA 路径: DSP 链上直接 fp4×fp8 MAC, 零解压",
     "节省 10-15% 带宽 + 延迟"),
    ("②", "MLA 硬件流水线",
     "GPU: 6 个 CUDA kernel + 显存往返\nFPGA: 6 级硬件流水, KV 不离开片上 SRAM",
     "消除 30 μs / 层 kernel launch 开销"),
    ("③", "硬件 KV Cache 寻址",
     "GPU: vLLM Block Table 软件管理, CPU 参与\nFPGA: {session,layer,seq} → 物理地址硬件生成",
     "零 CPU 参与, 零内存碎片"),
]
for i, (num, title, body, gain) in enumerate(cards):
    y = Inches(1.6 + i * 1.85)
    # Number circle
    add_box(s, Inches(0.5), y, Inches(0.8), Inches(0.8), ACCENT)
    add_text(s, Inches(0.5), y, Inches(0.8), Inches(0.8),
             num, size=32, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    # Title
    add_text(s, Inches(1.5), y, Inches(8.0), Inches(0.5),
             title, size=18, color=DARK_TEXT, bold=True)
    # Body
    tf = add_text(s, Inches(1.5), y + Inches(0.55), Inches(8.0), Inches(1.0),
                  "", size=12, color=MID_TEXT)
    for j, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(12); p.font.color.rgb = MID_TEXT
        p.font.name = FONT_CN; p.space_after = Pt(2)
    # Gain box
    add_box(s, Inches(9.7), y + Inches(0.1), Inches(3.1), Inches(1.5), LIGHT_BG)
    add_text(s, Inches(9.7), y + Inches(0.45), Inches(3.1), Inches(0.4),
             "收益", size=11, color=MID_TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.8), y + Inches(0.85), Inches(2.9), Inches(0.7),
             gain, size=11, color=GREEN_OK, bold=True,
             align=PP_ALIGN.CENTER)



# ═════════════════════════════════════════════════════════════
# Slide 11: DS V4 Architecture × FPGA Advantage
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 11, "DeepSeek V4 架构创新, FPGA 是天然加速器",
            "DS V4 的每一层创新都在减少 HBM/算力/内存压力 — 恰好是 FPGA vs GPU 的结构性优势")

add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.35),
         "10 大架构创新对照: GPU 能做到吗? FPGA 能做到吗?",
         size=13, color=DARK_TEXT, bold=True)

ds_rows = [
    ("fp4 权重 (E2M1)",      "解压→FP8 后计算 (~15% 开销)",  "DSP 原生 fp4×fp8 MAC, 零解压",               "✅ DONE"),
    ("MoE (top-6/384)",      "batch 大时 all-to-all 通信瓶颈", "32-chip pipeline + C2C dispatch/reduce",       "✅ DONE"),
    ("MLA (KV 压缩 56×)",    "6 kernel launch × 61 层 = 1.8ms",  "硬件 KV cache 寻址 + SRAM keeper",              "⚠️ partial"),
    ("DSA / CSA 稀疏注意力", "warp divergence, GPU 做不到",   "DSP 阵列早停电路 — partial sum < 阈值即跳过",  "🔥 P0 (0.5 人月)"),
    ("Engram (hash lookup)",  "GPU 只有 HBM, 无 LPDDR 可用",  "FPGA LPDDR+HBM+SRAM 三级存储, O(1) hash 硬件",  "🔥 P0 (1.0 人月)"),
    ("mHC (层间 highway)",    "每层多一次 kernel launch",     "嵌入 32-chip pipeline, 零额外延迟",             "🔥 P1 (0.5 人月)"),
    ("MTP (多 token 预测)",   "serial speculative decoding",  "并行 2-4 lm_head, 一次验证",                    "🔥 P1 (0.8 人月)"),
    ("KV Cache offload",     "GPU HBM 存不下 → recompute",   "SSD → PCIe DMA → HBM, <0.1s session 切换",     "P2 (0.5 人月)"),
]
add_table(s, Inches(0.5), Inches(1.95),
          [Inches(3.0), Inches(3.6), Inches(3.8), Inches(1.8)],
          ["DS V4 创新", "GPU 的局限", "FPGA 独占优势", "状态 / 投入"],
          ds_rows, header_size=11, row_size=10, row_h=Inches(0.5))

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
         "→ 结论: DS V4 的 fp4/MoE/MLA/DSA/Engram/mHC 六大创新, FPGA 比 GPU 更适合做硬件加速",
         size=13, color=ACCENT, bold=True)


# ═════════════════════════════════════════════════════════════
# Slide 12: FPGA 独占优势实现路线图
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 12, "FPGA 独占优势实现路线 (P0→P1→P2)",
            "总投入 4.5 人月, 3 阶段交付, 每阶段独立可验证")

# P0 block
add_box(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.6), RGBColor(0xFF, 0xF0, 0xE0))
add_text(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(0.35),
         "P0 — 上板前必须实现 (1.5 人月)", size=15, color=ACCENT, bold=True)
add_text(s, Inches(0.7), Inches(2.05), Inches(11.7), Inches(0.35),
         "DSA/CSA 稀疏注意力早停 | 插入 fp4_systolic_array 的 accumulate 阶段 | ~200 LUT + 1 DSP | 0.5 人月",
         size=10, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(2.35), Inches(11.7), Inches(0.35),
         "Engram O(1) hash lookup engine   | LPDDR + HBM + SRAM 三级存储 | ~500 LUT + 2 BRAM | 1.0 人月",
         size=10, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(2.65), Inches(11.7), Inches(0.35),
         "→ 这两个是 FPGA vs GPU 的结构性差异点, BP 核心论据. 实现后 context 从 128K→1M 时 compute 不增长.",
         size=10, color=ACCENT, bold=True)

# P1 block
add_box(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.4), RGBColor(0xE0, 0xF0, 0xFF))
add_text(s, Inches(0.7), Inches(3.45), Inches(11.9), Inches(0.35),
         "P1 — Phase 2 实现 (1.3 人月)", size=15, color=ACCENT2, bold=True)
add_text(s, Inches(0.7), Inches(3.85), Inches(11.7), Inches(0.35),
         "mHC layer mixer   | 嵌入 pipeline 零额外延迟, GPU 每层多一次 kernel launch | ~100 DSP | 0.5 人月",
         size=10, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(4.15), Inches(11.7), Inches(0.35),
         "MTP multi-token head | 并行 2-4 lm_head, 推理吞吐 1.5-2× | ~500 DSP | 0.8 人月",
         size=10, color=DARK_TEXT)
add_text(s, Inches(0.7), Inches(4.45), Inches(11.7), Inches(0.35),
         "→ FPGA 独占性能放大器. GPU 用 serial speculative decoding 达不到同样的加速比.",
         size=10, color=ACCENT2, bold=True)

# P2 block
add_box(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.2), RGBColor(0xEE, 0xEE, 0xEE))
add_text(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(0.35),
         "P2 — Phase 3 实现 (1.7 人月)", size=15, color=MID_TEXT, bold=True)
add_text(s, Inches(0.7), Inches(5.45), Inches(11.7), Inches(0.35),
         "MLA full pipeline (Q/K/V low-rank + RoPE) 1.2 人月  |  KV offload DMA engine 0.5 人月",
         size=10, color=DARK_TEXT)

# Summary box
add_box(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5), ACCENT)
add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         "4.5 person-months total -> 6 FPGA-exclusive advantages -> performance gap vs GPU from close to significant lead",
         size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════
# Slide 13: 实测数据 — 18 配置验证 (THE MONEY SLIDE)
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 13, "实测: 18 配置端到端仿真验证",
            "scripts/run_e2e_validation.py · 90s × 18 configs · 数据见 docs/e2e_validation_results.json")

# Headline metrics
metric_card(s, Inches(0.5), Inches(1.7), Inches(3.0), Inches(1.5),
            "×6.0", "Agent 负载 TPS 提升", "961 → 5,790 tok/s")
metric_card(s, Inches(3.7), Inches(1.7), Inches(3.0), Inches(1.5),
            "×318", "Burst TTFT P95 改善", "150.3s → 473 ms",
            value_color=GREEN_OK)
metric_card(s, Inches(6.9), Inches(1.7), Inches(3.0), Inches(1.5),
            "×2.9", "Agent Accept Rate", "24% → 70%")
metric_card(s, Inches(10.1), Inches(1.7), Inches(2.7), Inches(1.5),
            "0", "新增硬件成本", "纯软件层优化",
            value_color=ACCENT2)

# Detail table from e2e
def get(key, m, default='—'):
    v = E2E.get(key, {}).get(m)
    return v if v is not None else default

rows = [
    ("baseline",  f"{int(get('agent | baseline','output_tps',0))}",
                  f"{get('agent | baseline','accept_rate',0):.0f}%",
                  f"{get('agent | baseline','ttft_p95',0):.0f}",
                  f"{int(get('burst | baseline','ttft_p95',0))}"),
    ("+ KV 扩容 (D)", f"{int(get('agent | +D','output_tps',0))}",
                  f"{get('agent | +D','accept_rate',0):.0f}%",
                  f"{get('agent | +D','ttft_p95',0):.0f}",
                  f"{int(get('burst | +D','ttft_p95',0))}"),
    ("+ 微批 (C)", f"{int(get('agent | +D+C','output_tps',0))}",
                  f"{get('agent | +D+C','accept_rate',0):.0f}%",
                  f"{get('agent | +D+C','ttft_p95',0):.0f}",
                  f"{int(get('burst | +D+C','ttft_p95',0))}"),
    ("+ Hot Replication (A)", f"{int(get('agent | +D+C+A','output_tps',0))}",
                  f"{get('agent | +D+C+A','accept_rate',0):.0f}%",
                  f"{get('agent | +D+C+A','ttft_p95',0):.0f}",
                  f"{int(get('burst | +D+C+A','ttft_p95',0))}"),
    ("+ Pipeline Cloning ×2", f"{int(get('agent | +all+PC2','output_tps',0))}",
                  f"{get('agent | +all+PC2','accept_rate',0):.0f}%",
                  f"{get('agent | +all+PC2','ttft_p95',0):.0f}",
                  f"{int(get('burst | +all+PC2','ttft_p95',0))}"),
    ("+ Pipeline Cloning ×4", f"{int(get('agent | +all+PC4','output_tps',0))}",
                  f"{get('agent | +all+PC4','accept_rate',0):.0f}%",
                  f"{get('agent | +all+PC4','ttft_p95',0):.0f}",
                  f"{int(get('burst | +all+PC4','ttft_p95',0))}"),
]
add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.4),
         "累积优化堆栈实测 (Agent 4 req/s + Burst 20 req/s):",
         size=14, color=DARK_TEXT, bold=True)
add_table(s, Inches(0.5), Inches(3.95),
          [Inches(3.5), Inches(2.0), Inches(1.8), Inches(2.5), Inches(2.5)],
          ["配置", "Agent TPS", "Agent Accept", "Agent TTFT_p95", "Burst TTFT_p95"],
          rows, header_size=11, row_size=10, row_h=Inches(0.32))

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
         "注: 单 pipeline 峰值 17,445 tok/s; Pipeline Cloning ×N 把聚合峰值乘 N, burst+PC4 达 28,981 tok/s",
         size=10, color=MID_TEXT)


# ═════════════════════════════════════════════════════════════
# Slide 14: 优化堆栈逻辑
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 14, "四层软件优化, 零硬件改动",
            "每层都是独立可关的开关, 客户按需启用")

cards = [
    ("D", "KV Cache 扩容",
     "blocks/chip: 4,096 → 22,528 (×5.5)",
     "解锁 session 上限 ~16 → ~144",
     "改 1 个常量"),
    ("C", "调度地板移除",
     "MIN_DECODE_BATCH 4 → 1",
     "低并发场景立即流转, 不再饿死",
     "改 1 个常量"),
    ("A", "Hot Expert Replication",
     "Zipf α=1.0, Top-6 专家 ×8 副本",
     "K_pipeline 25.4 → 23.1 (-9%)",
     "调度器 + 权重布局 ~1 人月"),
    ("PC", "Pipeline Cloning",
     "32 chip 切 N 条独立 pipeline",
     "Prefill admission ×N, TTFT 大幅改善",
     "权重布局编译器 ~0.5 人月"),
]
for i, (code, title, what, why, cost) in enumerate(cards):
    y = Inches(1.7 + i * 1.25)
    # Code badge (smaller font for 2-letter "PC")
    add_box(s, Inches(0.5), y, Inches(1.2), Inches(1.05), ACCENT)
    badge_size = 36 if len(code) == 1 else 28
    add_text(s, Inches(0.5), y, Inches(1.2), Inches(1.05),
             code, size=badge_size, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    # Content
    add_text(s, Inches(1.9), y + Inches(0.05), Inches(7.0), Inches(0.4),
             title, size=15, color=DARK_TEXT, bold=True)
    add_text(s, Inches(1.9), y + Inches(0.45), Inches(7.0), Inches(0.35),
             what, size=11, color=MID_TEXT, font=FONT_EN)
    add_text(s, Inches(1.9), y + Inches(0.75), Inches(7.0), Inches(0.35),
             "→ " + why, size=11, color=GREEN_OK)
    # Cost
    add_box(s, Inches(9.2), y + Inches(0.15), Inches(3.6), Inches(0.75), LIGHT_BG)
    add_text(s, Inches(9.2), y + Inches(0.25), Inches(3.6), Inches(0.3),
             "实施成本", size=10, color=MID_TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.2), y + Inches(0.55), Inches(3.6), Inches(0.3),
             cost, size=11, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# Slide 15: $/M token 对比
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 15, "成本: 量产 100 套后 $1.30 / 百万 token, 低于 API",
            "纯硬件 TCO, 不含 IP 摊薄. 详见 docs/tco_per_million_tokens.md §5.2")

# Bar chart
chart_data = CategoryChartData()
chart_data.categories = [
    'NVIDIA H100\n(云租赁)',
    'Ascend 910C',
    'FPGA baseline\n(10 套)',
    'FPGA 优化\n(10 套)',
    'FPGA 优化\n(100 套)',
    'FPGA ASIC\n(终局)',
    'DeepSeek V4\n公开 API',
]
# Values: USD per million tokens
chart_data.add_series('$/百万 token', (16.0, 15.0, 6.0, 1.73, 1.30, 0.5, 1.46))

chart = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.7), Inches(8.0), Inches(5.0),
    chart_data
).chart
chart.has_legend = False
chart.has_title = False
# Format
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(11)
plot.data_labels.font.bold = True
plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
plot.gap_width = 60
# Axis
chart.value_axis.maximum_scale = 20.0
chart.value_axis.minimum_scale = 0.0
chart.value_axis.tick_labels.font.size = Pt(10)
chart.category_axis.tick_labels.font.size = Pt(10)

# Right side: cost composition
add_text(s, Inches(8.8), Inches(1.7), Inches(4.3), Inches(0.4),
         "成本拆解 (10 套量产, 单套年度)",
         size=14, color=DARK_TEXT, bold=True)
items = [
    ("硬件 BOM",     "¥1.46M / 套",         ACCENT),
    ("研发摊薄",     "¥450K / 套 (5 年)",  ACCENT2),
    ("年运营",       "¥273K (电+维)",       MID_TEXT),
    ("3 年硬件 TCO", "≈ $107K / 套",        DARK_TEXT),
    ("有效年产出",  "72B tokens (混合负载)", DARK_TEXT),
    ("$/M token",   "$1.30 (优化, 100 套)",     GREEN_OK),
    ("$/M token",   "$6.0 (baseline, 10 套)",     YELLOW),
]
for i, (k, v, col) in enumerate(items):
    y = Inches(2.25 + i * 0.55)
    add_text(s, Inches(8.8), y, Inches(2.0), Inches(0.4),
             k, size=11, color=MID_TEXT)
    add_text(s, Inches(10.8), y, Inches(2.3), Inches(0.4),
             v, size=11, color=col, bold=True, font=FONT_EN)


# ═════════════════════════════════════════════════════════════
# Slide 16: 市场
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 16, "目标市场: 1,500-5,000 套 / 年",
            "服务于 \"必须本地部署 + GPU 不可获取\" 的客户群")

# Bottom-up TAM
add_text(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
         "Bottom-Up: 按客户画像估算 (套数 / 年)",
         size=14, color=DARK_TEXT, bold=True)
rows = [
    ("中国科技公司出海 (TikTok/阿里云国际)",      "50-90",     "170-330",   "650-1,200"),
    ("央企/国企海外机构 (一带一路/四大行)",       "40-75",     "120-230",   "380-750"),
    ("金融/医疗/政府 私有部署",                    "30-60",     "100-200",   "400-800"),
    ("大模型公司自建集群 (DeepSeek/Moonshot)",     "20-40",     "80-150",    "300-600"),
    ("教育/科研机构",                                "5-10",      "40-100",    "140-300"),
    ("合计 (套数 / 年)",                             "145-275",   "510-1,010", "1,870-3,650"),
]
add_table(s, Inches(0.5), Inches(2.2),
          [Inches(4.5), Inches(2.3), Inches(2.3), Inches(2.3)],
          ["客户画像", "1-2 年", "3-5 年", "5-10 年"],
          rows, header_size=11, row_size=10, row_h=Inches(0.42))

# Right: hit rate
add_box(s, Inches(11.7), Inches(2.2), Inches(1.4), Inches(2.5), LIGHT_BG)
add_text(s, Inches(11.7), Inches(2.3), Inches(1.4), Inches(0.4),
         "需捕获",  size=11, color=MID_TEXT, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(11.7), Inches(2.65), Inches(1.4), Inches(0.6),
         "5-10%", size=22, color=ACCENT, bold=True,
         align=PP_ALIGN.CENTER, font=FONT_EN)
add_text(s, Inches(11.7), Inches(3.3), Inches(1.4), Inches(0.4),
         "= 100-500\n套 / 年",  size=10, color=MID_TEXT,
         align=PP_ALIGN.CENTER)

# Lower part: revenue projection
add_text(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
         "营收 / 利润测算 (按 100 套 / 年, 中期):",
         size=14, color=DARK_TEXT, bold=True)
items = [
    ("硬件营收 (¥1.70M / 套)",          "¥170M / 年"),
    ("IP license 摊薄收入",            "~¥10M / 年"),
    ("硬件成本 (¥1.18M / 套)",         "¥118M / 年"),
    ("毛利",                            "¥52M / 年 (31%)"),
    ("年运营总成本 (5 人核心团队)",     "¥15M / 年"),
    ("净利润",                          "¥37M / 年"),
]
for i, (k, v) in enumerate(items):
    col = 0 if i < 3 else 1
    row = i % 3
    x = Inches(0.5 + col * 6.5)
    y = Inches(5.8 + row * 0.42)
    add_text(s, x, y, Inches(4.0), Inches(0.4),
             k, size=11, color=MID_TEXT)
    add_text(s, x + Inches(4.0), y, Inches(2.0), Inches(0.4),
             v, size=11, color=DARK_TEXT, bold=True, font=FONT_EN,
             align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════
# Slide 17: 三阶段路线
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 17, "三阶段路线: FPGA → eASIC → ASIC",
            "分阶段释放风险, NRE 投入随市场验证递增")

phases = [
    ("Phase 1", "FPGA 原型",      "现在 - 12 月", "¥7M",   "5+ 种子客户",
     "10 套出货", "成本 $6.0/M",
     "RTL 验证 fp4+MLA, 真实负载校准, 低 NRE 试错"),
    ("Phase 2a", "FPGA 量产",     "12 - 24 月",    "¥3M",   "10-100 套出货",
     "成本 $1.30-1.73/M", "毛利 35-45%",
     "供应链建立, 软件栈成熟, 进入正循环"),
    ("Phase 2b", "eASIC 降本",   "18 - 30 月",    "¥35M",  "4 FPGA→1 eASIC",
     "成本 $0.6/M", "毛利 50%+",
     "结构化 ASIC, 复用 HBM2e+EMIB, NRE 仅金属层"),
    ("Phase 3", "全定制 ASIC",   "30 - 60 月",    "¥150M", "市场验证后",
     "成本 $0.3-0.5/M", "毛利 60%+",
     "HBM3@7nm, 终极成本, 需 1000+ 套订单支撑"),
]
for i, (phase, name, time, nre, milestone, perf1, perf2, desc) in enumerate(phases):
    y = Inches(1.7 + i * 1.35)
    # Phase badge
    bg = ACCENT if i == 0 else (ACCENT2 if i < 3 else MID_TEXT)
    add_box(s, Inches(0.5), y, Inches(1.6), Inches(1.15), bg)
    add_text(s, Inches(0.5), y + Inches(0.1), Inches(1.6), Inches(0.4),
             phase, size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, Inches(0.5), y + Inches(0.5), Inches(1.6), Inches(0.4),
             name, size=12, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), y + Inches(0.85), Inches(1.6), Inches(0.3),
             time, size=10, color=WHITE,
             align=PP_ALIGN.CENTER)
    # NRE box
    add_box(s, Inches(2.25), y + Inches(0.05), Inches(1.5), Inches(1.05), LIGHT_BG)
    add_text(s, Inches(2.25), y + Inches(0.15), Inches(1.5), Inches(0.3),
             "NRE", size=10, color=MID_TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.25), y + Inches(0.4), Inches(1.5), Inches(0.5),
             nre, size=22, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, Inches(2.25), y + Inches(0.85), Inches(1.5), Inches(0.3),
             milestone, size=8, color=MID_TEXT,
             align=PP_ALIGN.CENTER)
    # Description
    add_text(s, Inches(3.95), y + Inches(0.1), Inches(5.5), Inches(0.4),
             desc, size=12, color=DARK_TEXT, bold=True)
    add_text(s, Inches(3.95), y + Inches(0.55), Inches(5.5), Inches(0.3),
             perf1 + "  |  " + perf2, size=11, color=GREEN_OK)


# ═════════════════════════════════════════════════════════════
# Slide 18: 团队 + 工作量 + 实施完成度
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 18, "团队配置: 8 人 × 10 月 = 80 人月",
            "RTL + 软件并行推进, 已完成软件栈 ~50% (仿真验证)")

# Workload breakdown
add_text(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.4),
         "工作量拆分", size=14, color=DARK_TEXT, bold=True)
work = [
    ("RTL 工程师 × 5",     "50 人月",  "13 个 RTL 模块"),
    ("软件工程师 × 2",     "20 人月",  "调度器 + API + 优化栈"),
    ("权重布局编译器 × 1", "10 人月",  "fp4 量化 + 副本放置"),
    ("合计",               "80 人月",  "10 个月并行交付"),
]
for i, (role, pm, what) in enumerate(work):
    y = Inches(2.2 + i * 0.5)
    add_text(s, Inches(0.5), y, Inches(2.5), Inches(0.4),
             role, size=12, color=DARK_TEXT, bold=(i == 3))
    add_text(s, Inches(3.0), y, Inches(1.3), Inches(0.4),
             pm, size=12, color=ACCENT, bold=True, font=FONT_EN)
    add_text(s, Inches(4.3), y, Inches(2.5), Inches(0.4),
             what, size=10, color=MID_TEXT)

# Vertical divider between two columns
add_box(s, Inches(7.0), Inches(1.8), Emu(15000), Inches(4.6), RGBColor(0xDD, 0xDD, 0xDD))

# Software completion status (right)
add_text(s, Inches(7.3), Inches(1.7), Inches(5.6), Inches(0.4),
         "软件栈已完成度 (2026/05 仿真验证)",
         size=14, color=DARK_TEXT, bold=True)
done = [
    ("推理引擎 (libfpga.so 接口)",                       "✓"),
    ("Continuous Batching 调度器",                         "✓"),
    ("KV Cache 扩容 (§4.6.1 解法 D)",                      "✓"),
    ("调度地板移除 (§4.6.1 解法 C)",                       "✓"),
    ("Hot Expert Replication (§4.6.1 解法 A)",             "✓"),
    ("Pipeline Cloning (§4.8.x)",                          "✓"),
    ("Chip 0 admission 解析模型",                          "✓"),
    ("OpenAI API 兼容层",                                  "○"),
    ("LangChain / Dify 生态适配",                          "○"),
    ("权重布局编译器 (待 RTL 上板后)",                     "○"),
]
for i, (item, status) in enumerate(done):
    y = Inches(2.2 + i * 0.4)
    col = GREEN_OK if status == "✓" else YELLOW
    add_text(s, Inches(7.3), y, Inches(0.4), Inches(0.35),
             status, size=14, color=col, bold=True, font=FONT_EN)
    add_text(s, Inches(7.8), y, Inches(5.0), Inches(0.35),
             item, size=11, color=DARK_TEXT)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "→ 软件栈核心优化已仿真验证完成. RTL 启动 Phase 1 后立即可上板对接.",
         size=12, color=ACCENT, bold=True)


# ═════════════════════════════════════════════════════════════
# Slide 19: 风险 + 对冲
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 19, "三大风险与对冲方案")

risks = [
    ("①", "fp4 Python 仿真已 PASS, 尚需上板",
     "group=16, alpha=1.0, cosine=0.99554; 真实 DSP rounding 未验证",
     "Phase 1 实验 1: 1 卡 1 层上板对比. 备选: group=16 保持, 敏感层退回 fp8",
     RED_BAD),
    ("②", "RTL 50 人月可行性",
     "13 个模块 + 5 人 10 月, 业界经验偏紧 30-50%",
     "Phase 1 单卡跑通 6 模块 + Signal Tap 加速调试. 备选: 增 1 人月 (¥1M)",
     YELLOW),
    ("③", "模型架构演进",
     "DeepSeek V5/V6 可能改 MoE Top-K 或 attention 变体",
     "FPGA Partial Reconfig 数十 ms 升级, 比 ASIC 流片 (18 月) 灵活 100×",
     GREEN_OK),
]
for i, (num, title, what, mitigation, color) in enumerate(risks):
    y = Inches(1.7 + i * 1.7)
    # Number circle
    add_box(s, Inches(0.5), y, Inches(0.7), Inches(0.7), color)
    add_text(s, Inches(0.5), y, Inches(0.7), Inches(0.7),
             num, size=28, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    # Title
    add_text(s, Inches(1.4), y, Inches(11.0), Inches(0.4),
             title, size=15, color=DARK_TEXT, bold=True)
    # What
    add_text(s, Inches(1.4), y + Inches(0.45), Inches(11.0), Inches(0.4),
             "风险: " + what, size=11, color=MID_TEXT)
    # Mitigation
    add_text(s, Inches(1.4), y + Inches(0.85), Inches(11.0), Inches(0.6),
             "对冲: " + mitigation, size=11, color=GREEN_OK)


# ═════════════════════════════════════════════════════════════
# Slide 20: 里程碑
# ═════════════════════════════════════════════════════════════
s = add_slide()
title_block(s, 20, "10 个月路线图 + 3 个 Go/No-Go 闸门")

# Timeline bar
bar_y = Inches(2.0)
bar_h = Inches(0.5)
bar_w = Inches(12.0)
add_box(s, Inches(0.7), bar_y, bar_w, bar_h, LIGHT_BG)

phases = [
    ("M1-2", "单卡验证", 0.20,  ACCENT, "G1: fp4 精度"),
    ("M3-4", "8 卡节点",  0.40,  ACCENT2, "G2: HBM 带宽"),
    ("M5-6", "双节点",    0.60,  ACCENT2, ""),
    ("M7-8", "全 32 卡",  0.80,  ACCENT2, "G3: 端到端 SLA"),
    ("M9-10","生产化",    1.0,   GREEN_OK, "首批客户交付"),
]
prev_x = 0.7
for label, name, end_frac, color, gate in phases:
    end_x = 0.7 + 12.0 * end_frac
    width = end_x - prev_x
    add_box(s, Inches(prev_x), bar_y, Inches(width), bar_h, color)
    add_text(s, Inches(prev_x), bar_y + Inches(0.1), Inches(width), Inches(0.3),
             name, size=12, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(prev_x), bar_y + Inches(0.55), Inches(width), Inches(0.3),
             label, size=10, color=MID_TEXT, align=PP_ALIGN.CENTER, font=FONT_EN)
    if gate:
        # Gate marker
        x_gate = Inches(end_x - 0.6)
        add_text(s, x_gate, bar_y + Inches(0.95), Inches(1.2), Inches(0.4),
                 "▼", size=14, color=RED_BAD, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, x_gate, bar_y + Inches(1.3), Inches(1.2), Inches(0.4),
                 gate, size=10, color=RED_BAD, bold=True,
                 align=PP_ALIGN.CENTER)
    prev_x = end_x

# Gate details below
add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
         "Go / No-Go 闸门 (任一不达标 → 重新评估方案):",
         size=14, color=DARK_TEXT, bold=True)
gates = [
    ("G1 (M2)", "fp4 精度", "Python 已达 0.99554; 上板 cosine ≥ 0.995",
     "失败 → 部分层退回 fp8 (吞吐 -5%) 或终止"),
    ("G2 (M4)", "HBM 带宽", "AGM 039 HBM 有效带宽 ≥ 80% × 920 GB/s = 736 GB/s",
     "失败 → 重新核算 decode 性能, 可能需要降低并发"),
    ("G3 (M8)", "端到端 SLA", "TTFT P95 ≤ 500 ms, agent 场景 accept ≥ 80%",
     "失败 → 调度器调优 1-2 月, 不动 RTL"),
]
for i, (g, what, criteria, fallback) in enumerate(gates):
    y = Inches(5.05 + i * 0.65)
    add_box(s, Inches(0.5), y, Inches(1.2), Inches(0.55), RED_BAD)
    add_text(s, Inches(0.5), y, Inches(1.2), Inches(0.55),
             g, size=11, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    add_text(s, Inches(1.85), y, Inches(1.5), Inches(0.55),
             what, size=12, color=DARK_TEXT, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.4), y, Inches(5.5), Inches(0.55),
             criteria, size=10, color=MID_TEXT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.0), y, Inches(4.0), Inches(0.55),
             fallback, size=10, color=YELLOW,
             anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════
# Slide 21: Ask
# ═════════════════════════════════════════════════════════════
s = add_slide(DARK_BG)
add_box(s, 0, 0, prs.slide_width, Inches(0.15), ACCENT)
# Title
add_text(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6),
         "Ask", size=18, color=ACCENT, bold=True)
add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0),
         "一期 ¥7M, 10 个月, 3 个验收闸门",
         size=40, color=WHITE, bold=True)

# Three columns: 钱怎么花 / 产出 / 退出
sections = [
    ("Phase 1 资金分配",
     [("¥3.3M", "RTL IP 开发 (5 人 × 10 月)"),
      ("¥1.5M", "软件 / 驱动 / 调度 (2 人 × 10 月)"),
      ("¥0.7M", "开发板 + Quartus License + 工具"),
      ("¥0.65M","权重布局编译器"),
      ("¥0.85M","测试 / 验证 / 杂费")]),
    ("Phase 1 产出",
     [("32 颗", "AGM 039 上板 8 卡 × 4 = 32 芯片"),
      ("13 个", "RTL IP 模块 (可复用至 eASIC)"),
      ("5 家",  "种子客户 POC 验证"),
      ("10 套", "首批量产硬件出货"),
      ("$1.0-6.0", "/百万 token 区间 (优化-baseline)")]),
    ("退出路径",
     [("Phase 2", "FPGA 量产, 营收 ¥170M / 年 @100 套"),
      ("Phase 2b","eASIC 降本, 毛利 50%+ 进入二级市场"),
      ("Phase 3", "ASIC 流片, $50M 营收 / 年 + 行业并购"),
      ("IP 出售", "RTL IP 转售/授权 ¥20-50M"),
      ("最坏",     "¥10-15M 沉没, RTL IP 可转其他加速场景")]),
]
for i, (title, items) in enumerate(sections):
    x = Inches(0.5 + i * 4.27)
    # Card
    add_box(s, x, Inches(2.6), Inches(4.1), Inches(4.0),
            RGBColor(0x22, 0x32, 0x52))
    # Title
    add_text(s, x, Inches(2.7), Inches(4.1), Inches(0.45),
             title, size=15, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER)
    # Top divider
    add_box(s, x + Inches(0.4), Inches(3.2), Inches(3.3), Emu(15000), ACCENT)
    # Items
    for j, (val, lbl) in enumerate(items):
        y = Inches(3.4 + j * 0.55)
        if val:
            add_text(s, x + Inches(0.3), y, Inches(1.3), Inches(0.4),
                     val, size=14, color=ACCENT, bold=True,
                     align=PP_ALIGN.LEFT, font=FONT_EN)
            add_text(s, x + Inches(1.6), y, Inches(2.4), Inches(0.4),
                     lbl, size=10, color=WHITE,
                     anchor=MSO_ANCHOR.MIDDLE)

# Bottom: contact
add_text(s, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4),
         "联系: liyan@huahuan.com  ·  下一步: 5 月底完成开发板采购, 6 月启动 Phase 1",
         size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════
# Save
# ═════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), '..', 'docs',
                            'fpga_inference_bp_v2.pptx')
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
