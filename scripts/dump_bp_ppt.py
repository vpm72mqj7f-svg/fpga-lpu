"""Dump every slide's text/shapes/tables for visual inspection without opening PPT."""
from pptx import Presentation
from pptx.util import Emu

prs = Presentation('docs/fpga_inference_bp.pptx')
SW = prs.slide_width / 914400  # inches
SH = prs.slide_height / 914400

print(f"Deck: {len(prs.slides)} slides, {SW:.1f}\" x {SH:.1f}\"")
print("=" * 80)

for idx, slide in enumerate(prs.slides, 1):
    print(f"\n┌─ SLIDE {idx:>2} " + "─" * 70)
    # Collect text by approximate Y position (top to bottom)
    items = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            top_in = shape.top / 914400 if shape.top else 0
            left_in = shape.left / 914400 if shape.left else 0
            txt = shape.text_frame.text.replace('\n', ' | ').strip()
            # Get first run font size & color for context
            try:
                p = shape.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else None
                fs = run.font.size.pt if run and run.font.size else None
            except Exception:
                fs = None
            items.append((top_in, left_in, fs, txt))
        elif shape.has_table:
            tbl = shape.table
            top_in = shape.top / 914400 if shape.top else 0
            cells = []
            for row in tbl.rows:
                row_txt = [c.text_frame.text.strip().replace('\n', ' ') for c in row.cells]
                cells.append(row_txt)
            items.append((top_in, 0, None, f"<TABLE {len(cells)}x{len(cells[0])}>"))
            for r_idx, row in enumerate(cells):
                items.append((top_in + 0.001*r_idx, 0.5, None,
                              "  | " + " | ".join(row)))
        elif shape.has_chart:
            top_in = shape.top / 914400 if shape.top else 0
            chart = shape.chart
            cats = list(chart.plots[0].categories)
            ser_vals = [list(s.values) for s in chart.series]
            items.append((top_in, 0, None,
                          f"<CHART cats={cats}>"))
            for s_idx, vals in enumerate(ser_vals):
                items.append((top_in + 0.001 + 0.001*s_idx, 0.5, None,
                              f"  series {s_idx}: {vals}"))

    items.sort(key=lambda x: (round(x[0], 1), x[1]))
    for top, left, fs, txt in items:
        fs_str = f"{fs:.0f}pt" if fs else "    "
        if len(txt) > 95:
            txt = txt[:92] + "..."
        print(f"│ y={top:4.1f} {fs_str:>6} │ {txt}")
print("\n" + "=" * 80)
